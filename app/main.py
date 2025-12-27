# app/main.py
from __future__ import annotations

import os
import json
import shutil
import asyncio
from typing import Any, Dict, List, Optional

from fastapi import (
    FastAPI,
    Depends,
    File,
    Form,
    HTTPException,
    UploadFile,
)
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import FileResponse, StreamingResponse
from sqlalchemy.orm import Session
from dotenv import load_dotenv

from app.core.config import settings
from app.db.database import SessionLocal, engine, Base
from app.db import crud, models
from app.ai.ai_manager import AIManager
from app.ai import tools as ai_tools
from app.ai.mcp_client import mcp_client, MCPClient
from app.utils.logger import logger, log_api_call, chat_logger
from app.utils.context_manager import ContextManager

# OCR 功能(延迟导入,避免启动时加载)
def get_ocr_module():
    try:
        from app.utils.ocr import ocr_image, is_ocr_available
        return ocr_image, is_ocr_available
    except ImportError:
        return None, lambda: False

load_dotenv()

# 数据库初始化(可选)
# 通过环境变量控制是否自动初始化数据库
AUTO_INIT_DB = os.getenv("AUTO_INIT_DB", "0") == "1"
if AUTO_INIT_DB:
    Base.metadata.create_all(bind=engine)

app = FastAPI(title=settings.APP_NAME, version=settings.APP_VERSION)

# CORS
app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],  # 如需限制,可改为具体域名
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

def get_db():
    db = SessionLocal()
    try:
        yield db
    finally:
        db.close()

def parse_bool(value: Optional[Any]) -> Optional[bool]:
    if value is None:
        return None
    if isinstance(value, bool):
        return value
    if isinstance(value, str):
        return value.strip().lower() in {"1", "true", "yes", "on"}
    return bool(value)

ai_manager = AIManager()

# MCP 服务器启动事件
@app.on_event("startup")
async def startup_event():
    """应用启动时加载 MCP 服务器配置(不自动启动,等待前端按需启动)"""
    try:
        db = SessionLocal()
        saved_config = crud.get_setting(db, "mcp_servers")
        if saved_config:
            servers_config = json.loads(saved_config.value)
            for config in servers_config:
                if config.get("enabled", True):
                    name = config.get("name", "")
                    command = config.get("command", "")
                    args = config.get("args", [])
                    env = config.get("env", {})
                    if name and command:
                        # 只添加配置,不启动服务器
                        mcp_client.add_server(name, command, args, env)
                        print(f"[MCP] 服务器 {name} 配置已加载")
        db.close()
    except Exception as e:
        chat_logger.error(f"[MCP] 加载配置失败: {e}")

@app.on_event("shutdown")
async def shutdown_event():
    """应用关闭时停止所有 MCP 服务"""
    await mcp_client.stop_all()

# ========== 基础接口 ==========

@app.post("/init-database")
def init_database():
    """手动初始化数据库(创建所有表)"""
    try:
        Base.metadata.create_all(bind=engine)
        return {"success": True, "message": "数据库初始化成功"}
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"数据库初始化失败: {str(e)}")

@app.get("/models")
def get_models():
    return {
        "default": settings.AI_MODEL,
        "models": settings.ai_models,
    }

@app.get("/provider/config")
def get_provider_config():
    """
    原有接口:返回当前运行时的 provider 配置.
    现在的实现:仅返回 .env 的静态配置是否存在,真正的多 Provider 管理在 /providers 系列接口.
    """
    has_key = bool(settings.AI_API_KEY.strip())
    return {
        "api_base": settings.AI_API_BASE,
        "has_key": has_key,
        "default_model": settings.AI_MODEL,
        "models": settings.ai_models,
    }

@app.post("/provider/config")
def set_provider_config(
    api_base: Optional[str] = Form(None),
    api_key: Optional[str] = Form(None),
    model: Optional[str] = Form(None),
):
    """
    原有接口:运行时覆盖当前 provider.
    为了兼容旧前端,这里仍然支持 set_provider,但实际上更推荐使用 /providers 管理.
    """
    try:
        ai_manager.set_provider(
            api_base=api_base or settings.AI_API_BASE,
            api_key=api_key or settings.AI_API_KEY,
            default_model=model or settings.AI_MODEL,
        )
    except ValueError as e:
        raise HTTPException(status_code=400, detail=str(e))

    return {"success": True}

# ========== 会话管理 ==========

@app.post("/conversations")
def create_conversation(
    title: Optional[str] = Form("新对话"),
    model: Optional[str] = Form(None),
    db: Session = Depends(get_db),
):
    """
    创建新对话:
    - 若最新对话没有任何消息,则复用该对话,避免无限新增空对话.
    - 否则创建新对话.
    返回字段:
    - conversation: 对话信息
    - reused: 是否复用了现有空对话
    """
    latest = crud.get_latest_conversation(db)
    if latest:
        msg_count = crud.get_conversation_message_count(db, latest.id)
        if msg_count == 0:
            # 复用最新的空对话
            if title and title != latest.title:
                latest = crud.update_conversation_title(db, latest.id, title)
            if model and model != latest.model:
                latest = crud.update_conversation_model(db, latest.id, model)
            return {"conversation": latest.to_dict(), "reused": True}

    conv = crud.create_conversation(db, title=title, model=model)
    return {"conversation": conv.to_dict(), "reused": False}

@app.get("/conversations")
def list_conversations(db: Session = Depends(get_db)):
    conversations = crud.get_conversations(db)
    return [conv.to_dict() for conv in conversations]

@app.delete("/conversations/{conversation_id}")
def delete_conversation(conversation_id: int, db: Session = Depends(get_db)):
    crud.delete_conversation(db, conversation_id)
    return {"success": True}

@app.put("/conversations/{conversation_id}")
def update_conversation(
    conversation_id: int,
    title: Optional[str] = Form(None),
    model: Optional[str] = Form(None),
    is_pinned: Optional[bool] = Form(None),
    enable_knowledge_base: Optional[bool] = Form(None),
    enable_mcp: Optional[bool] = Form(None),
    enable_web_search: Optional[bool] = Form(None),
    provider_id: Optional[int] = Form(None),
    db: Session = Depends(get_db),
):
    conv = crud.get_conversation(db, conversation_id)
    if not conv:
        raise HTTPException(status_code=404, detail="Conversation not found")

    if title is not None:
        conv.title = title
    if model is not None:
        conv.model = model
    if is_pinned is not None:
        conv.is_pinned = parse_bool(is_pinned) or False
    if enable_knowledge_base is not None:
        conv.enable_knowledge_base = parse_bool(enable_knowledge_base) or False
    if enable_mcp is not None:
        conv.enable_mcp = parse_bool(enable_mcp) or False
    if enable_web_search is not None:
        conv.enable_web_search = parse_bool(enable_web_search) or False

    if provider_id is not None:
        conv.provider_id = provider_id

    db.commit()
    db.refresh(conv)
    return conv.to_dict()

@app.post("/conversations/{conversation_id}/title")
def update_conversation_title(
    conversation_id: int,
    title: str = Form(...),
    db: Session = Depends(get_db),
):
    conv = crud.update_conversation_title(db, conversation_id, title)
    if not conv:
        raise HTTPException(status_code=404, detail="Conversation not found")
    return conv.to_dict()

@app.post("/conversations/{conversation_id}/auto-title")
def auto_generate_conversation_title(
    conversation_id: int,
    model: Optional[str] = Form(None),
    first_user_message: Optional[str] = Form(None),
    db: Session = Depends(get_db),
):
    """
    基于对话内容自动生成标题
    """
    import traceback
    
    conversation = crud.get_conversation(db, conversation_id)
    if not conversation:
        raise HTTPException(status_code=404, detail="Conversation not found")
    
    # 获取对话的前几条消息用于生成标题
    messages_db = crud.get_messages(db, conversation_id)
    # 如果数据库中消息太少,但前端传来的 first_user_message,则把它作为最小上下文
    if len(messages_db) < 2 and not first_user_message:
        raise HTTPException(status_code=400, detail="对话内容不足,无法生成标题")

    # 取前4条消息作为上下文
    context_messages = messages_db[:4]
    if len(context_messages) < 2 and first_user_message:
        # 创建一个临时的用户消息对象样式用于生成标题
        class _TmpMsg:
            def __init__(self, role, content):
                self.role = role
                self.content = content

        tmp = _TmpMsg('user', first_user_message)
        # 把临时消息插到最前面
        context_messages = [tmp] + context_messages
    context_text = "\n".join([f"{m.role}: {m.content[:200]}" for m in context_messages])  # 限制每条消息长度
    
    # 构建标题生成的提示
    title_prompt = f"""请为以下对话生成一个简洁的标题(不超过10个字):

{context_text}

要求:
1. 标题要准确概括对话主题
2. 使用中文
3. 不超过10个字
4. 不要包含引号或特殊符号
5. 直接返回标题,不要其他内容

标题:"""

    try:
        # 配置AI Provider
        _configure_ai_provider_for_conversation(db, conversation)
        
        # 确定使用的模型:优先参数,其次设置中的 auto_title_model,再次会话全局默认
        selected_model = model
        if not selected_model:
            setting_model = crud.get_setting(db, "auto_title_model")
            if setting_model and setting_model.value and setting_model.value != "current":
                selected_model = setting_model.value
        if not selected_model:
            selected_model = conversation.model or settings.AI_MODEL

        use_model = selected_model
        
        # 调用AI生成标题
        title_messages = [{"role": "user", "content": title_prompt}]
        
        # 添加超时和重试机制
        try:
            result = ai_manager.chat(title_messages, model=use_model, stream=False)
            # 从结果中提取内容
            if isinstance(result, dict) and "content" in result:
                generated_title = result["content"]
            else:
                generated_title = str(result)
        except Exception as api_error:
            chat_logger.warning(f"标题生成API调用失败: {api_error}")
            # 如果API调用失败,使用简单的标题生成逻辑
            user_message = context_messages[0].content if context_messages else "新对话"
            generated_title = user_message[:15] + "..." if len(user_message) > 15 else user_message
        
        # 清理生成的标题
        generated_title = str(generated_title).strip().replace('"', '').replace("'", "").replace("标题：", "").replace("标题:", "")
        if len(generated_title) > 10:
            generated_title = generated_title[:10]
        
        # 确保标题不为空
        if not generated_title or generated_title.isspace():
            generated_title = "新对话"
        
        # 更新对话标题
        conv = crud.update_conversation_title(db, conversation_id, generated_title)
        return {"title": generated_title, "conversation": conv.to_dict()}
        
    except Exception as e:
        chat_logger.error(f"标题生成错误: {e}")
        raise HTTPException(status_code=500, detail=f"标题生成失败: {str(e)}")

@app.post("/conversations/{conversation_id}/pin")
def update_conversation_pin(
    conversation_id: int,
    is_pinned: bool = Form(...),
    db: Session = Depends(get_db),
):
    is_pinned = bool(parse_bool(is_pinned))
    conv = crud.update_conversation_pin(db, conversation_id, is_pinned)

    if not conv:
        raise HTTPException(status_code=404, detail="Conversation not found")
    return conv.to_dict()

@app.post("/conversations/{conversation_id}/model")
def update_conversation_model(
    conversation_id: int,
    model: Optional[str] = Form(None),
    db: Session = Depends(get_db),
):
    conv = crud.update_conversation_model(db, conversation_id, model)
    if not conv:
        raise HTTPException(status_code=404, detail="Conversation not found")
    return conv.to_dict()

# 新增:更新会话功能开关(知识库/MCP/联网搜索)
@app.post("/conversations/{conversation_id}/features")
def update_conversation_features(
    conversation_id: int,
    enable_knowledge_base: Optional[bool] = Form(None),
    enable_mcp: Optional[bool] = Form(None),
    enable_web_search: Optional[bool] = Form(None),
    db: Session = Depends(get_db),
):
    enable_knowledge_base = parse_bool(enable_knowledge_base)
    enable_mcp = parse_bool(enable_mcp)
    enable_web_search = parse_bool(enable_web_search)

    conv = crud.update_conversation_features(
        db,
        conversation_id,
        enable_knowledge_base=enable_knowledge_base,
        enable_mcp=enable_mcp,
        enable_web_search=enable_web_search,
    )

    if not conv:
        raise HTTPException(status_code=404, detail="Conversation not found")
    return conv.to_dict()

# 新增:绑定 Provider 到会话
@app.post("/conversations/{conversation_id}/provider")
def set_conversation_provider(
    conversation_id: int,
    provider_id: Optional[str] = Form(None),
    db: Session = Depends(get_db),
):
    parsed_provider: Optional[int]
    if provider_id in (None, "", "null", "undefined"):
        parsed_provider = None
    else:
        try:
            parsed_provider = int(provider_id)
        except Exception:
            raise HTTPException(status_code=400, detail="Invalid provider_id")

    conv = crud.set_conversation_provider(db, conversation_id, parsed_provider)
    if not conv:
        raise HTTPException(status_code=404, detail="Conversation not found")
    return conv.to_dict()

# ========== 消息与聊天 ==========

@app.get("/conversations/{conversation_id}/messages")
def get_messages(conversation_id: int, db: Session = Depends(get_db)):
    messages = crud.get_messages(db, conversation_id)
    return [msg.to_dict() for msg in messages]

@app.post("/conversations/{conversation_id}/messages/partial")
def save_partial_message(
    conversation_id: int,
    content: str = Form(...),
    model: Optional[str] = Form(None),
    thinking_content: Optional[str] = Form(None),
    db: Session = Depends(get_db),
):
    """
    保存部分消息(用于用户中断流式输出时保存已生成的内容)
    """
    conversation = crud.get_conversation(db, conversation_id)
    if not conversation:
        raise HTTPException(status_code=404, detail="Conversation not found")
    
    if not content or not content.strip():
        return {"success": False, "message": "内容为空,不保存"}
    
    # 添加中断标记
    content_with_mark = content.strip() + "\n\n[输出被中断]"
    
    token_info = {
        "model": model or "unknown",
        "input_tokens": 0,
        "output_tokens": 0,
        "total_tokens": 0
    }
    
    msg = crud.create_message(
        db, conversation_id, "assistant", content_with_mark, token_info,
        tool_calls=None, thinking_content=thinking_content
    )
    
    return {"success": True, "message_id": msg.id}

def _configure_ai_provider_for_conversation(
    db: Session,
    conversation: models.Conversation,
    override_provider_id: Optional[int] = None,
) -> None:
    """
    根据会话绑定的 provider 或覆盖参数,配置 AIManager 当前使用的 provider.
    """
    provider: Optional[models.Provider] = None

    if override_provider_id is not None:
        provider = crud.get_provider(db, override_provider_id)
    elif conversation.provider_id:
        provider = crud.get_provider(db, conversation.provider_id)

    # 如果没有找到 provider,尝试使用第一个可用的 provider
    if not provider:
        all_providers = crud.list_providers(db)
        if all_providers:
            provider = all_providers[0]

    if provider:
        ai_manager.set_provider(
            api_base=provider.api_base,
            api_key=provider.api_key,
            default_model=conversation.model or provider.default_model,
        )
    else:
        # 使用全局默认
        if not settings.AI_API_BASE:
            raise HTTPException(status_code=400, detail="未配置任何 Provider,请先在设置中添加 Provider")
        ai_manager.set_provider(
            api_base=settings.AI_API_BASE,
            api_key=settings.AI_API_KEY,
            default_model=conversation.model or settings.AI_MODEL,
        )

def _execute_chat_with_tools(
    messages: List[Dict[str, Any]], 
    tools_list: List[Dict[str, Any]], 
    model: Optional[str],
    conversation_id: int,
    db: Session
) -> tuple[str, Dict[str, Any]]:
    """
    执行带工具的对话,包括工具调用循环
    """
    content, token_info, _ = _execute_chat_with_tools_streaming(messages, tools_list, model, conversation_id, db)
    return content, token_info

def _execute_chat_with_tools_streaming(
    messages: List[Dict[str, Any]], 
    tools_list: List[Dict[str, Any]], 
    model: Optional[str],
    conversation_id: int,
    db: Session,
    on_tool_call=None
) -> tuple[str, Dict[str, Any], List[Dict[str, Any]]]:
    """
    执行带工具的对话,包括工具调用循环,返回工具调用信息
    """
    import json
    from app.ai import tools as ai_tools
    
    # 累计token统计
    total_input_tokens = 0
    total_output_tokens = 0
    
    # 工具调用信息
    tool_calls_info = []
    
    # 工具调用循环
    current_messages = messages.copy()
    max_iterations = 5  # 防止无限循环
    
    for iteration in range(max_iterations):
        # 调用模型
        data = ai_manager.run_with_tools(current_messages, tools=tools_list, model=model, stream=False)
        
        # 累计token统计
        usage = data.get("usage", {})
        total_input_tokens += usage.get("prompt_tokens", 0)
        total_output_tokens += usage.get("completion_tokens", 0)
        
        message = data["choices"][0]["message"]
        # 清理消息格式，只保留必要字段
        clean_message = {
            "role": "assistant",
            "content": message.get("content") or ""
        }
        if message.get("tool_calls"):
            clean_message["tool_calls"] = message.get("tool_calls")
        current_messages.append(clean_message)
        
        # 检查是否有工具调用
        tool_calls = message.get("tool_calls")
        if not tool_calls:
            # 没有工具调用,返回最终结果
            final_content = message.get("content", "")
            token_info = {
                "model": model or "default",
                "input_tokens": total_input_tokens,
                "output_tokens": total_output_tokens,
                "total_tokens": total_input_tokens + total_output_tokens
            }
            return final_content, token_info, tool_calls_info
        
        # 执行工具调用
        for tool_call in tool_calls:
            function_name = tool_call["function"]["name"]
            function_args = json.loads(tool_call["function"]["arguments"])
            
            # 记录工具调用信息
            tool_info = {
                "name": function_name,
                "args": function_args,
                "status": "running"
            }
            tool_calls_info.append(tool_info)
            
            # 执行工具
            try:
                result = _execute_tool(function_name, function_args, conversation_id, db)
                tool_info["status"] = "success"
                tool_info["result_preview"] = result[:100] + "..." if len(result) > 100 else result
            except Exception as e:
                result = f"工具执行失败: {str(e)}"
                tool_info["status"] = "error"
                tool_info["error"] = str(e)
            
            # 添加工具调用结果到消息历史
            current_messages.append({
                "role": "tool",
                "tool_call_id": tool_call["id"],
                "content": result
            })
    
    # 如果达到最大迭代次数,返回最后的消息
    final_content = current_messages[-1].get("content", "达到最大工具调用次数限制")
    token_info = {
        "model": model or "default",
        "input_tokens": total_input_tokens,
        "output_tokens": total_output_tokens,
        "total_tokens": total_input_tokens + total_output_tokens
    }
    return final_content, token_info, tool_calls_info

def _execute_tool(function_name: str, function_args: Dict[str, Any], conversation_id: int, db: Session) -> str:
    """
    执行具体的工具调用
    """
    from app.ai import tools as ai_tools
    
    try:
        # 检查是否是 MCP 工具调用(格式:mcp_服务器名_工具名)
        if function_name.startswith("mcp_"):
            # 使用 mcp_client 的解析方法来正确处理包含中文的服务器名
            server_name, tool_name = mcp_client.parse_tool_name(function_name)
            
            if server_name and tool_name:
                
                # 异步调用 MCP 工具
                import asyncio
                try:
                    loop = asyncio.get_event_loop()
                except RuntimeError:
                    loop = asyncio.new_event_loop()
                    asyncio.set_event_loop(loop)
                
                result = loop.run_until_complete(
                    mcp_client.call_tool(server_name, tool_name, function_args)
                )
                
                if result.get("error"):
                    return f"MCP 工具执行失败: {result['error']}"
                
                # 格式化结果
                mcp_result = result.get("result", {})
                if isinstance(mcp_result, dict):
                    content = mcp_result.get("content", [])
                    if content:
                        # 提取文本内容
                        texts = []
                        for item in content:
                            if isinstance(item, dict) and item.get("type") == "text":
                                texts.append(item.get("text", ""))
                            elif isinstance(item, str):
                                texts.append(item)
                        return "\n".join(texts) if texts else json.dumps(mcp_result, ensure_ascii=False, indent=2)
                    return json.dumps(mcp_result, ensure_ascii=False, indent=2)
                return str(mcp_result)
            else:
                return f"无效的 MCP 工具名称格式: {function_name}"
        
        if function_name == "get_local_time":
            return ai_tools.run_get_local_time_tool()
        
        elif function_name == "calculate_expression":
            expression = function_args.get("expression", "")
            return ai_tools.run_calculator_tool(expression)
        
        elif function_name == "search_knowledge":
            query = function_args.get("query", "")
            kb_id = function_args.get("kb_id")
            top_k = function_args.get("top_k", 5)
            
            chat_logger.info(f"知识库搜索: query={query}, kb_id={kb_id}")
            
            # 获取知识库使用的 embedding 模型
            embedding_model = None
            embedding_provider = None
            
            # 1. 先尝试从知识库文档中获取 embedding 模型
            docs = crud.list_knowledge_documents(db, kb_id=kb_id) if kb_id else crud.list_knowledge_documents(db)
            if docs:
                for doc in docs:
                    if doc.embedding_model:
                        embedding_model = doc.embedding_model
                        break
            
            # 2. 查找 Provider 中的 embedding 模型
            all_providers = crud.list_providers(db)
            
            for provider in all_providers:
                if provider.models_config:
                    try:
                        import json as _json
                        config = _json.loads(provider.models_config)
                        for model_name in config.keys():
                            # 检查是否是 embedding 模型
                            if "embedding" in model_name.lower() or "embed" in model_name.lower():
                                if not embedding_model:
                                    embedding_model = model_name
                                embedding_provider = provider
                                break
                    except Exception:
                        pass
                if embedding_provider:
                    break
            
            # 3. 如果还是没找到 embedding Provider,使用第一个 Provider
            if not embedding_provider and all_providers:
                embedding_provider = all_providers[0]
            
            if not embedding_model:
                return "未配置向量模型,无法进行知识库搜索。请在 Provider 设置中添加 embedding 模型(如 text-embedding-3-small)。"
            
            # 检查知识库是否有向量数据
            chunks_count = db.query(models.KnowledgeChunk).count()
            if chunks_count == 0:
                return "知识库中没有向量数据。请重新上传文档,并在上传时选择向量模型(如 text-embedding-3-small)。"
            
            if embedding_provider:
                ai_manager.set_provider(
                    api_base=embedding_provider.api_base,
                    api_key=embedding_provider.api_key,
                    default_model=embedding_provider.default_model,
                )
            else:
                return "未配置任何 Provider,无法进行知识库搜索"
            
            # 创建embedding函数
            final_embedding_model = embedding_model
            def embedding_fn(texts):
                try:
                    return ai_manager.create_embedding(texts, model=final_embedding_model)
                except Exception as e:
                    chat_logger.error(f"Embedding调用失败: {e}")
                    return None
            
            return ai_tools.run_search_knowledge_tool(
                query=query,
                kb_id=kb_id,
                top_k=top_k,
                embedding_fn=embedding_fn,
                use_graph=False  # 已移除知识图谱功能
            )
        
        elif function_name == "web_search":
            query = function_args.get("query", "")
            source = function_args.get("source", "duckduckgo")
            return ai_tools.run_web_search_tool(query=query, source=source)
        
        else:
            return f"未知的工具: {function_name}"
            
    except Exception as e:
        return f"工具执行错误: {str(e)}"

def _build_tools_for_conversation(
    conversation: models.Conversation,
    enable_knowledge_base: Optional[bool],
    enable_mcp: Optional[bool],
    enable_web_search: Optional[bool],
) -> List[Dict[str, Any]]:
    """
    根据会话默认开关 + 本次请求参数,决定启用哪些 tools.
    优先使用本次请求参数,如果为 None 则回退到 conversation 的设置.
    """
    kb_flag = (
        enable_knowledge_base
        if enable_knowledge_base is not None
        else conversation.enable_knowledge_base
    )
    mcp_flag = (
        enable_mcp if enable_mcp is not None else conversation.enable_mcp
    )
    web_flag = (
        enable_web_search
        if enable_web_search is not None
        else conversation.enable_web_search
    )

    # 获取 MCP 工具列表
    mcp_tools = mcp_client.get_all_tools() if mcp_flag else None

    return ai_tools.get_tools(
        enable_knowledge_base=kb_flag,
        enable_mcp=mcp_flag,
        enable_web_search=web_flag,
        mcp_tools=mcp_tools,
    )

def _get_conversation_files_context(
    db: Session, 
    conversation_id: int,
    current_model: Optional[str] = None,
    model_supports_vision: bool = False
) -> tuple[str, List[Dict[str, Any]], List[Dict[str, Any]]]:
    """
    读取对话关联的文件内容,返回格式化的上下文字符串、图片列表和需要视觉识别的文件列表.
    
    返回:
        - file_context: 文本文件的内容
        - image_files: 图片文件列表 [{"filepath": ..., "filename": ...}]
        - files_need_vision: 需要视觉识别的文件列表(PDF/Word/PPT等本地解析无内容的)
    """
    from app.utils.document_parser import extract_text_from_file
    
    files = crud.get_uploaded_files(db, conversation_id)
    if not files:
        return "", [], []
    
    file_contents = []
    image_files = []
    files_need_vision = []  # 需要视觉识别的文件(PDF/Word/PPT)
    total_length = 0
    max_total_length = 50000  # 限制总长度约 50K 字符
    max_per_file = 15000  # 每个文件最多 15K 字符
    
    # 图片扩展名
    image_extensions = {'.png', '.jpg', '.jpeg', '.gif', '.bmp', '.webp'}
    # 支持视觉识别的文档扩展名
    vision_doc_extensions = {'.pdf', '.doc', '.docx', '.ppt', '.pptx'}
    
    for file_record in files:
        if total_length >= max_total_length:
            break
            
        try:
            if not os.path.exists(file_record.filepath):
                continue
            
            # 检查是否是图片文件
            ext = os.path.splitext(file_record.filename)[1].lower()
            if ext in image_extensions:
                image_files.append({
                    "filepath": file_record.filepath,
                    "filename": file_record.filename
                })
                continue
            
            # 尝试提取文本内容
            content = extract_text_from_file(file_record.filepath, extract_images=False)
            
            # 检查是否是支持视觉识别的文档且没有提取到内容
            if ext in vision_doc_extensions and (not content or not content.strip()):
                # 文档没有文本内容,需要视觉识别
                files_need_vision.append({
                    "filepath": file_record.filepath,
                    "filename": file_record.filename,
                    "file_type": ext[1:]  # 去掉点号:pdf, docx, pptx 等
                })
                continue
            
            if not content or not content.strip():
                continue
            
            # 截断过长的内容
            if len(content) > max_per_file:
                content = content[:max_per_file] + "\n...(内容已截断)"
            
            file_contents.append(f"【文件: {file_record.filename}】\n{content}")
            total_length += len(content)
            
        except Exception as e:
            chat_logger.warning(f"读取文件 {file_record.filename} 失败: {e}")
            continue
    
    text_context = "\n\n".join(file_contents) if file_contents else ""
    return text_context, image_files, files_need_vision

def _recognize_images_with_ocr(
    image_files: List[Dict[str, Any]],
    use_ocr: bool = True
) -> Tuple[str, List[Dict[str, Any]]]:
    """
    使用本地 OCR 识别图片中的文字
    
    Args:
        image_files: 图片文件列表
        use_ocr: 是否启用 OCR
        
    Returns:
        (ocr_results, remaining_files): OCR 识别结果和剩余需要视觉模型处理的文件
    """
    if not use_ocr or not image_files:
        return "", image_files
    
    ocr_image, is_ocr_available = get_ocr_module()
    if not is_ocr_available():
        return "", image_files
    
    ocr_results = []
    remaining_files = []
    
    for img_info in image_files:
        filepath = img_info["filepath"]
        filename = img_info["filename"]
        
        try:
            text = ocr_image(filepath)
            if text and text.strip():  # 有内容就用
                ocr_results.append(f"【图片: {filename}】\n{text}")
            else:
                # OCR 没有识别到文字,交给视觉模型
                remaining_files.append(img_info)
        except Exception:
            remaining_files.append(img_info)
    
    return "\n\n".join(ocr_results), remaining_files

def _recognize_images_with_ocr_stream(
    image_files: List[Dict[str, Any]],
    use_ocr: bool = True
):
    """
    使用本地 OCR 识别图片中的文字(流式版本)
    
    Yields:
        事件字典: type 可以是 start, progress, chunk, result, end
    """
    if not use_ocr or not image_files:
        return
    
    ocr_image, is_ocr_available = get_ocr_module()
    if not is_ocr_available():
        return
    
    total = len(image_files)
    yield {"type": "start", "model": "本地OCR", "total": total, "file_type": "image"}
    
    for idx, img_info in enumerate(image_files):
        filepath = img_info["filepath"]
        filename = img_info["filename"]
        
        yield {"type": "progress", "message": f"正在OCR识别 ({idx + 1}/{total}): {filename}"}
        
        try:
            text = ocr_image(filepath)
            if text and text.strip():
                # 分块输出
                for line in text.split('\n'):
                    if line.strip():
                        yield {"type": "chunk", "content": line + "\n"}
                yield {"type": "result", "content": f"【图片: {filename}】\n{text}", "has_text": len(text.strip()) > 20}
            else:
                yield {"type": "result", "content": "", "has_text": False}
        except Exception:
            yield {"type": "result", "content": "", "has_text": False}
    
    yield {"type": "end"}

def _recognize_docs_with_ocr(
    doc_files: List[Dict[str, Any]]
) -> str:
    """
    使用OCR识别文档内容(同步版本)
    将文档转为图片后进行OCR
    """
    results = []
    for event in _recognize_docs_with_ocr_stream(doc_files):
        if event["type"] == "result" and event.get("content"):
            results.append(event["content"])
    return "\n\n".join(results) if results else ""

def _recognize_docs_with_ocr_stream(
    doc_files: List[Dict[str, Any]]
):
    """
    使用OCR识别文档内容(流式版本)
    将文档转为图片后进行OCR
    """
    if not doc_files:
        return
    
    ocr_image, is_ocr_available = get_ocr_module()
    if not is_ocr_available():
        yield {"type": "error", "message": "OCR模块不可用"}
        return
    
    total_files = len(doc_files)
    yield {"type": "start", "model": "本地OCR", "total": total_files, "file_type": "document"}
    
    for file_idx, doc_info in enumerate(doc_files):
        try:
            filepath = doc_info["filepath"]
            filename = doc_info["filename"]
            ext = os.path.splitext(filename)[1].lower()
            
            yield {"type": "progress", "message": f"正在OCR识别文档 ({file_idx + 1}/{total_files}): {filename}"}
            
            try:
                from PIL import Image
                import io
                import tempfile
                
                images = []
                
                # 根据文件类型转换为图片
                if ext == '.pdf':
                    from pdf2image import convert_from_path
                    import platform
                    
                    poppler_path = None
                    if platform.system() == "Windows":
                        poppler_path = r"C:\poppler\poppler-24.08.0\Library\bin"
                        if not os.path.exists(poppler_path):
                            poppler_path = None
                    
                    images = convert_from_path(filepath, first_page=1, last_page=10, dpi=150, poppler_path=poppler_path)
                    
                elif ext in ['.ppt', '.pptx']:
                    images = _convert_ppt_to_images(filepath)
                    
                elif ext in ['.doc', '.docx']:
                    images = _convert_word_to_images(filepath)
                
                if not images:
                    yield {"type": "result", "content": f"【文档: {filename}】\n(无法转换为图片进行OCR识别)"}
                    continue
                
                # 对每页图片进行OCR
                page_contents = []
                for page_idx, img in enumerate(images):
                    yield {"type": "progress", "message": f"正在OCR识别 {filename} 第 {page_idx + 1} 页"}
                    
                    # 保存图片到临时文件
                    with tempfile.NamedTemporaryFile(suffix='.png', delete=False) as tmp:
                        img.save(tmp.name, format='PNG')
                        tmp_path = tmp.name
                    
                    try:
                        text = ocr_image(tmp_path)
                        if text and text.strip():
                            page_contents.append(f"[第 {page_idx + 1} 页]\n{text}")
                            # 分块输出
                            for line in text.split('\n'):
                                if line.strip():
                                    yield {"type": "chunk", "content": line + "\n"}
                    finally:
                        # 删除临时文件
                        try:
                            os.unlink(tmp_path)
                        except:
                            pass
                
                if page_contents:
                    result_content = f"【文档: {filename}(OCR识别)】\n" + "\n\n".join(page_contents)
                    yield {"type": "result", "content": result_content}
                else:
                    yield {"type": "result", "content": f"【文档: {filename}】\n(OCR未识别到文字)"}
                    
            except ImportError as e:
                chat_logger.warning(f"文档OCR依赖未安装: {e}")
                yield {"type": "result", "content": f"【文档: {filename}】\n(缺少必要依赖: {str(e)})"}
            except Exception as e:
                chat_logger.warning(f"文档OCR失败: {e}")
                yield {"type": "result", "content": f"【文档: {filename}】\n(OCR识别失败: {str(e)})"}
                
        except Exception as e:
            chat_logger.warning(f"文档OCR处理失败: {e}")
            yield {"type": "result", "content": f"【文档: {doc_info.get('filename', '未知')}】\n(处理失败)"}
    
    yield {"type": "end"}

def _recognize_images_with_vision_model(
    db: Session,
    image_files: List[Dict[str, Any]],
    vision_model: str
) -> str:
    """
    使用视觉模型识别图片内容(同步版本)
    """
    results = []
    for event in _recognize_images_with_vision_model_stream(db, image_files, vision_model):
        if event["type"] == "result":
            results.append(event["content"])
    return "\n\n".join(results) if results else ""

def _recognize_images_with_vision_model_stream(
    db: Session,
    image_files: List[Dict[str, Any]],
    vision_model: str
):
    """
    使用视觉模型识别图片内容(流式版本,yield 进度事件)
    事件类型:
    - start: 开始识别
    - progress: 识别进度(流式内容)
    - result: 单个图片识别完成
    - end: 全部完成
    """
    import base64
    
    if not image_files or not vision_model:
        return
    
    total = len(image_files)
    yield {"type": "start", "model": vision_model, "total": total, "file_type": "image"}
    
    for idx, img_info in enumerate(image_files):
        try:
            filepath = img_info["filepath"]
            filename = img_info["filename"]
            
            yield {"type": "progress", "message": f"正在识别图片 ({idx + 1}/{total}): {filename}"}
            
            # 读取图片并转为 base64
            with open(filepath, "rb") as f:
                image_data = base64.b64encode(f.read()).decode("utf-8")
            
            # 获取图片 MIME 类型
            ext = os.path.splitext(filename)[1].lower()
            mime_map = {
                '.png': 'image/png',
                '.jpg': 'image/jpeg',
                '.jpeg': 'image/jpeg',
                '.gif': 'image/gif',
                '.bmp': 'image/bmp',
                '.webp': 'image/webp'
            }
            mime_type = mime_map.get(ext, 'image/png')
            
            # 构建视觉模型请求
            messages = [
                {
                    "role": "user",
                    "content": [
                        {
                            "type": "text",
                            "text": f"请详细描述这张图片的内容。图片文件名: {filename}"
                        },
                        {
                            "type": "image_url",
                            "image_url": {
                                "url": f"data:{mime_type};base64,{image_data}"
                            }
                        }
                    ]
                }
            ]
            
            # 调用视觉模型(流式)
            content_parts = []
            for chunk in ai_manager.chat(messages, model=vision_model, stream=True):
                if isinstance(chunk, dict):
                    chunk_content = chunk.get("content", "")
                else:
                    chunk_content = chunk
                if chunk_content:
                    content_parts.append(chunk_content)
                    yield {"type": "chunk", "content": chunk_content}
            
            content = "".join(content_parts)
            if content:
                yield {"type": "result", "content": f"【图片: {filename}】\n{content}"}
                
        except Exception as e:
            chat_logger.warning(f"识别图片 {img_info.get('filename', '未知')} 失败: {e}")
            yield {"type": "result", "content": f"【图片: {img_info.get('filename', '未知')}】\n(图片识别失败)"}
    
    yield {"type": "end"}

def _recognize_pdf_with_vision_model(
    db: Session,
    pdf_files: List[Dict[str, Any]],
    vision_model: str
) -> str:
    """
    使用视觉模型识别文档内容(同步版本)
    支持 PDF、Word、PPT
    """
    results = []
    for event in _recognize_docs_with_vision_model_stream(db, pdf_files, vision_model):
        if event["type"] == "result":
            results.append(event["content"])
    return "\n\n".join(results) if results else ""

def _recognize_docs_with_vision_model_stream(
    db: Session,
    doc_files: List[Dict[str, Any]],
    vision_model: str
):
    """
    使用视觉模型识别文档内容(流式版本)
    支持 PDF、Word (.doc/.docx)、PPT (.ppt/.pptx)
    每两页上下拼接成一张图片，这样可以识别更多页面
    """
    import base64
    
    if not doc_files or not vision_model:
        return
    
    total_files = len(doc_files)
    yield {"type": "start", "model": vision_model, "total": total_files, "file_type": "document"}
    
    for file_idx, doc_info in enumerate(doc_files):
        try:
            filepath = doc_info["filepath"]
            filename = doc_info["filename"]
            ext = os.path.splitext(filename)[1].lower()
            file_type = doc_info.get("file_type", ext[1:])
            
            yield {"type": "progress", "message": f"正在处理文档 ({file_idx + 1}/{total_files}): {filename}"}
            
            try:
                from PIL import Image
                import io
                
                images = []
                
                # 根据文件类型转换为图片
                if ext == '.pdf':
                    from pdf2image import convert_from_path
                    import platform
                    
                    # Windows 需要指定 poppler 路径
                    poppler_path = None
                    if platform.system() == "Windows":
                        poppler_path = r"C:\poppler\poppler-24.08.0\Library\bin"
                        if not os.path.exists(poppler_path):
                            poppler_path = None  # 让 pdf2image 尝试从 PATH 查找
                    
                    images = convert_from_path(filepath, first_page=1, last_page=10, dpi=150, poppler_path=poppler_path)
                    
                elif ext in ['.ppt', '.pptx']:
                    # PPT 转图片
                    images = _convert_ppt_to_images(filepath)
                    
                elif ext in ['.doc', '.docx']:
                    # Word 转图片
                    images = _convert_word_to_images(filepath)
                
                if not images:
                    yield {"type": "result", "content": f"【文档: {filename}】\n(无法转换为图片进行识别)"}
                    continue
                
                # 将图片两两拼接
                merged_images = []
                page_ranges = []
                
                for i in range(0, len(images), 2):
                    if i + 1 < len(images):
                        # 有两页,上下拼接
                        img1 = images[i]
                        img2 = images[i + 1]
                        
                        # 确保宽度一致(取较大的宽度)
                        max_width = max(img1.width, img2.width)
                        
                        # 如果宽度不一致,调整图片
                        if img1.width != max_width:
                            ratio = max_width / img1.width
                            img1 = img1.resize((max_width, int(img1.height * ratio)), Image.Resampling.LANCZOS)
                        if img2.width != max_width:
                            ratio = max_width / img2.width
                            img2 = img2.resize((max_width, int(img2.height * ratio)), Image.Resampling.LANCZOS)
                        
                        # 创建拼接后的图片
                        total_height = img1.height + img2.height + 20  # 20px 间隔
                        merged = Image.new('RGB', (max_width, total_height), 'white')
                        merged.paste(img1, (0, 0))
                        merged.paste(img2, (0, img1.height + 20))
                        
                        merged_images.append(merged)
                        page_ranges.append(f"{i + 1}-{i + 2}")
                    else:
                        # 只有一页,直接使用
                        merged_images.append(images[i])
                        page_ranges.append(f"{i + 1}")
                
                page_contents = []
                total_pages = len(merged_images)
                for idx, (merged_img, page_range) in enumerate(zip(merged_images, page_ranges)):
                    yield {"type": "progress", "message": f"正在识别 {filename} 第 {page_range} 页 ({idx + 1}/{total_pages})"}
                    
                    # 将图片转为 base64
                    buffer = io.BytesIO()
                    merged_img.save(buffer, format='PNG', optimize=True)
                    image_data = base64.b64encode(buffer.getvalue()).decode("utf-8")
                    
                    # 构建视觉模型请求
                    if "-" in page_range:
                        prompt_text = f"请识别并提取这张图片中的所有文字内容。这是 {filename} 的第 {page_range} 页(上下两页拼接在一起)。请按顺序输出页面中的文字,保持原有格式。"
                    else:
                        prompt_text = f"请识别并提取这个文档页面中的所有文字内容。这是 {filename} 的第 {page_range} 页。请直接输出页面中的文字,保持原有格式。"
                    
                    messages = [
                        {
                            "role": "user",
                            "content": [
                                {
                                    "type": "text",
                                    "text": prompt_text
                                },
                                {
                                    "type": "image_url",
                                    "image_url": {
                                        "url": f"data:image/png;base64,{image_data}"
                                    }
                                }
                            ]
                        }
                    ]
                    
                    # 调用视觉模型(流式)
                    content_parts = []
                    for chunk in ai_manager.chat(messages, model=vision_model, stream=True):
                        if isinstance(chunk, dict):
                            chunk_content = chunk.get("content", "")
                        else:
                            chunk_content = chunk
                        if chunk_content:
                            content_parts.append(chunk_content)
                            yield {"type": "chunk", "content": chunk_content}
                    
                    content = "".join(content_parts)
                    if content:
                        page_contents.append(f"[第 {page_range} 页]\n{content}")
                
                if page_contents:
                    result_content = f"【文档: {filename}(视觉识别)】\n" + "\n\n".join(page_contents)
                    yield {"type": "result", "content": result_content}
                    
            except ImportError as e:
                chat_logger.warning(f"文档视觉识别依赖未安装: {e}")
                yield {"type": "result", "content": f"【文档: {filename}】\n(缺少必要依赖: {str(e)})"}
            except Exception as e:
                chat_logger.warning(f"文档转图片失败: {e}")
                yield {"type": "result", "content": f"【文档: {filename}】\n(视觉识别失败: {str(e)})"}
                
        except Exception as e:
            chat_logger.warning(f"文档视觉识别失败: {e}")
            yield {"type": "result", "content": f"【文档: {doc_info.get('filename', '未知')}】\n(识别失败)"}
    
    yield {"type": "end"}

def _convert_ppt_to_images(filepath: str) -> List:
    """
    将 PPT/PPTX 转换为图片列表
    使用 python-pptx 渲染幻灯片为图片
    """
    from PIL import Image, ImageDraw, ImageFont
    import io
    
    images = []
    ext = os.path.splitext(filepath)[1].lower()
    
    if ext not in ['.ppt', '.pptx']:
        return images
    
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.enum.shapes import MSO_SHAPE_TYPE
        
        prs = Presentation(filepath)
        slide_width = prs.slide_width.pt if prs.slide_width else 960
        slide_height = prs.slide_height.pt if prs.slide_height else 540
        
        # 设置输出图片尺寸(按比例缩放)
        scale = 2  # 放大倍数以提高清晰度
        img_width = int(slide_width * scale)
        img_height = int(slide_height * scale)
        
        for slide_idx, slide in enumerate(prs.slides[:10]):  # 最多10页
            # 创建白色背景图片
            img = Image.new('RGB', (img_width, img_height), 'white')
            draw = ImageDraw.Draw(img)
            
            # 尝试加载字体
            try:
                font = ImageFont.truetype("arial.ttf", int(24 * scale))
                font_small = ImageFont.truetype("arial.ttf", int(16 * scale))
            except:
                try:
                    font = ImageFont.truetype("msyh.ttc", int(24 * scale))
                    font_small = ImageFont.truetype("msyh.ttc", int(16 * scale))
                except:
                    font = ImageFont.load_default()
                    font_small = font
            
            # 提取幻灯片中的文本和图片
            y_offset = int(50 * scale)
            
            for shape in slide.shapes:
                try:
                    # 处理图片
                    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                        try:
                            image_stream = io.BytesIO(shape.image.blob)
                            shape_img = Image.open(image_stream)
                            
                            # 计算位置和大小
                            left = int(shape.left.pt * scale) if shape.left else 0
                            top = int(shape.top.pt * scale) if shape.top else y_offset
                            width = int(shape.width.pt * scale) if shape.width else 200
                            height = int(shape.height.pt * scale) if shape.height else 150
                            
                            # 调整图片大小
                            shape_img = shape_img.resize((width, height), Image.Resampling.LANCZOS)
                            if shape_img.mode == 'RGBA':
                                img.paste(shape_img, (left, top), shape_img)
                            else:
                                img.paste(shape_img, (left, top))
                        except Exception:
                            pass
                    
                    # 处理文本框
                    if shape.has_text_frame:
                        text = shape.text_frame.text.strip()
                        if text:
                            left = int(shape.left.pt * scale) if shape.left else int(50 * scale)
                            top = int(shape.top.pt * scale) if shape.top else y_offset
                            
                            # 绘制文本
                            draw.text((left, top), text, fill='black', font=font_small)
                            y_offset = top + int(30 * scale)
                            
                except Exception:
                    continue
            
            images.append(img)
            
    except Exception as e:
        chat_logger.warning(f"PPT转图片失败: {e}")
    
    return images

def _convert_word_to_images(filepath: str) -> List:
    """
    将 Word 文档转换为图片列表
    使用 python-docx 提取内容并渲染为图片
    """
    from PIL import Image, ImageDraw, ImageFont
    import io
    
    images = []
    ext = os.path.splitext(filepath)[1].lower()
    
    if ext not in ['.doc', '.docx']:
        return images
    
    try:
        from docx import Document
        from docx.shared import Inches
        
        doc = Document(filepath)
        
        # 设置页面尺寸(A4 比例)
        page_width = 1600
        page_height = 2200
        margin = 100
        
        # 尝试加载字体
        try:
            font = ImageFont.truetype("arial.ttf", 28)
            font_title = ImageFont.truetype("arial.ttf", 36)
        except:
            try:
                font = ImageFont.truetype("msyh.ttc", 28)
                font_title = ImageFont.truetype("msyh.ttc", 36)
            except:
                font = ImageFont.load_default()
                font_title = font
        
        # 创建第一页
        current_img = Image.new('RGB', (page_width, page_height), 'white')
        draw = ImageDraw.Draw(current_img)
        y_offset = margin
        page_count = 0
        
        for element in doc.element.body:
            # 检查是否需要新页面
            if y_offset > page_height - margin * 2:
                images.append(current_img)
                page_count += 1
                if page_count >= 10:  # 最多10页
                    break
                current_img = Image.new('RGB', (page_width, page_height), 'white')
                draw = ImageDraw.Draw(current_img)
                y_offset = margin
            
            # 处理段落
            if element.tag.endswith('p'):
                for para in doc.paragraphs:
                    if para._element == element:
                        text = para.text.strip()
                        if text:
                            # 自动换行
                            max_width = page_width - margin * 2
                            lines = []
                            current_line = ""
                            
                            for char in text:
                                test_line = current_line + char
                                bbox = draw.textbbox((0, 0), test_line, font=font)
                                if bbox[2] - bbox[0] > max_width:
                                    if current_line:
                                        lines.append(current_line)
                                    current_line = char
                                else:
                                    current_line = test_line
                            if current_line:
                                lines.append(current_line)
                            
                            for line in lines:
                                if y_offset > page_height - margin:
                                    images.append(current_img)
                                    page_count += 1
                                    if page_count >= 10:
                                        break
                                    current_img = Image.new('RGB', (page_width, page_height), 'white')
                                    draw = ImageDraw.Draw(current_img)
                                    y_offset = margin
                                
                                draw.text((margin, y_offset), line, fill='black', font=font)
                                y_offset += 40
                            
                            y_offset += 20  # 段落间距
                        break
            
            # 处理图片
            if element.tag.endswith('drawing') or element.tag.endswith('pict'):
                for rel in doc.part.rels.values():
                    if "image" in rel.reltype:
                        try:
                            image_data = rel.target_part.blob
                            shape_img = Image.open(io.BytesIO(image_data))
                            
                            # 调整图片大小
                            max_img_width = page_width - margin * 2
                            max_img_height = 400
                            shape_img.thumbnail((max_img_width, max_img_height), Image.Resampling.LANCZOS)
                            
                            if y_offset + shape_img.height > page_height - margin:
                                images.append(current_img)
                                page_count += 1
                                if page_count >= 10:
                                    break
                                current_img = Image.new('RGB', (page_width, page_height), 'white')
                                draw = ImageDraw.Draw(current_img)
                                y_offset = margin
                            
                            if shape_img.mode == 'RGBA':
                                current_img.paste(shape_img, (margin, y_offset), shape_img)
                            else:
                                current_img.paste(shape_img, (margin, y_offset))
                            y_offset += shape_img.height + 20
                        except Exception:
                            pass
        
        # 添加最后一页
        if y_offset > margin:
            images.append(current_img)
            
    except Exception as e:
        chat_logger.warning(f"Word转图片失败: {e}")
    
    return images
    
    yield {"type": "end"}

@app.post("/conversations/{conversation_id}/chat")
@log_api_call
def chat_with_conversation(
    conversation_id: int,
    user_text: str = Form(...),
    model: Optional[str] = Form(None),

    # 新增:本次请求的功能开关(可覆盖会话默认)
    enable_knowledge_base: Optional[bool] = Form(None),
    enable_mcp: Optional[bool] = Form(None),
    enable_web_search: Optional[bool] = Form(None),
    web_search_source: Optional[str] = Form(None),  # 搜索源
    
    # 新增:深度思考开关
    enable_thinking: Optional[bool] = Form(None),
    
    # 新增:强制启用视觉识别(用于PDF等文件)
    force_vision_recognition: Optional[bool] = Form(None),

    # 新增:指定本次使用的 provider(可选)
    provider_id: Optional[int] = Form(None),

    # 新增:是否流式输出(默认 False)
    stream: bool = Form(False),

    db: Session = Depends(get_db),
):
    from datetime import datetime
    start_time = datetime.now()

    enable_knowledge_base = parse_bool(enable_knowledge_base)
    enable_mcp = parse_bool(enable_mcp)
    enable_web_search = parse_bool(enable_web_search)
    enable_thinking = parse_bool(enable_thinking) or False
    force_vision_recognition = parse_bool(force_vision_recognition) or False
    stream = parse_bool(stream) or False
    
    # 记录聊天请求
    tools_enabled = {
        "knowledge_base": enable_knowledge_base,
        "mcp": enable_mcp,
        "web_search": enable_web_search,
        "thinking": enable_thinking
    }

    logger.log_chat_request(conversation_id, user_text, model, tools_enabled)
    
    conversation = crud.get_conversation(db, conversation_id)
    if not conversation:
        logger.log_error(Exception("Conversation not found"), f"对话ID {conversation_id} 不存在")
        raise HTTPException(status_code=404, detail="Conversation not found")

    # 1. 写入用户消息
    try:
        user_msg = crud.create_message(db, conversation_id, "user", user_text)
        logger.log_database_operation("CREATE", "messages", user_msg.id, {
            "role": "user", 
            "content_length": len(user_text)
        })
    except Exception as e:
        logger.log_error(e, "创建用户消息失败")
        raise

    # 2. 配置 Provider
    try:
        _configure_ai_provider_for_conversation(db, conversation, override_provider_id=provider_id)
        logger.log_performance("配置Provider", (datetime.now() - start_time).total_seconds())
    except Exception as e:
        logger.log_error(e, "配置Provider失败")
        raise

    # 3. 准备上下文消息(只包含完整问答对)
    context_messages = crud.get_context_messages(db, conversation_id)
    messages: List[Dict[str, Any]] = [
        {"role": m.role, "content": m.content} for m in context_messages
    ]
    
    # 3.5 读取对话关联的文件内容和图片
    # 获取当前使用的模型
    current_model = model or conversation.model or settings.AI_MODEL
    
    # 检查当前模型是否支持视觉
    model_supports_vision = False
    all_providers = crud.list_providers(db)
    for provider in all_providers:
        if provider.models_config:
            try:
                import json as _json
                config = _json.loads(provider.models_config)
                if current_model in config:
                    caps = config[current_model]
                    model_supports_vision = caps.get("vision", False)
                    break
            except:
                pass
    
    # 读取文件内容、图片列表和需要视觉识别的文档
    file_context, image_files, files_need_vision = _get_conversation_files_context(db, conversation_id)
    
    # 获取默认视觉模型(格式可能是 "provider_id:model_name" 或旧格式 "model_name")
    default_vision_model_setting = crud.get_setting(db, "default_vision_model")
    default_vision_model = None
    vision_provider_id = None
    if default_vision_model_setting and default_vision_model_setting.value:
        vision_value = default_vision_model_setting.value
        if ":" in vision_value:
            # 新格式:provider_id:model_name
            parts = vision_value.split(":", 1)
            try:
                vision_provider_id = int(parts[0])
                default_vision_model = parts[1]
            except ValueError:
                # 解析失败,当作旧格式处理
                default_vision_model = vision_value
        else:
            # 旧格式:只有 model_name,需要查找对应的 provider
            default_vision_model = vision_value
            # 遍历所有 provider 查找包含该模型的 provider
            for p in all_providers:
                if p.models_config:
                    try:
                        import json as _json
                        config = _json.loads(p.models_config)
                        if default_vision_model in config:
                            vision_provider_id = p.id
                            break
                    except:
                        pass
    
    # 准备需要视觉识别的文件列表(延迟到流式处理中执行)
    images_need_vision = []  # 需要视觉模型识别的图片
    docs_need_vision = []    # 需要视觉模型识别的文档(PDF/Word/PPT)
    docs_need_ocr = []       # 需要OCR识别的文档(没有文字的文档)
    
    # 处理图片文件
    if image_files:
        if model_supports_vision:
            # 当前模型支持视觉,直接在消息中包含图片(后续处理)
            pass  # 图片将在构建消息时处理
        else:
            # 当前模型不支持视觉
            if force_vision_recognition and default_vision_model:
                # 用户强制启用视觉识别，图片发给视觉模型
                images_need_vision = image_files
            # 否则默认使用 OCR(在后续处理中执行)
    
    # 处理文档文件
    if model_supports_vision:
        # 场景1：模型支持视觉
        # 没有文字的文档(files_need_vision)需要转图片直接发给AI
        if files_need_vision:
            docs_need_vision = files_need_vision
    else:
        # 模型不支持视觉
        if force_vision_recognition and default_vision_model:
            # 场景3：勾选了眼睛按钮，对所有文档使用视觉模型
            vision_doc_extensions = {'.pdf', '.doc', '.docx', '.ppt', '.pptx'}
            all_doc_files = []
            files = crud.get_uploaded_files(db, conversation_id)
            for file_record in files:
                ext = os.path.splitext(file_record.filename)[1].lower()
                if ext in vision_doc_extensions and os.path.exists(file_record.filepath):
                    all_doc_files.append({
                        "filepath": file_record.filepath,
                        "filename": file_record.filename,
                        "file_type": ext[1:]
                    })
            if all_doc_files:
                docs_need_vision = all_doc_files
        else:
            # 场景2：未勾选眼睛按钮，没有文字的文档用OCR
            if files_need_vision:
                docs_need_ocr = files_need_vision
    
    # 非流式模式下的处理
    image_context = ""
    doc_context = ""
    if not stream:
        # 处理图片
        if image_files and not model_supports_vision:
            if force_vision_recognition and default_vision_model:
                # 用户强制启用视觉识别
                if vision_provider_id:
                    vision_provider = crud.get_provider(db, vision_provider_id)
                    if vision_provider:
                        ai_manager.set_provider(
                            api_base=vision_provider.api_base,
                            api_key=vision_provider.api_key,
                            default_model=default_vision_model
                        )
                image_context = _recognize_images_with_vision_model(db, image_files, default_vision_model)
                if vision_provider_id and provider:
                    ai_manager.set_provider(
                        api_base=provider.api_base,
                        api_key=provider.api_key,
                        default_model=model
                    )
            else:
                # 默认使用 OCR
                ocr_context, _ = _recognize_images_with_ocr(image_files, use_ocr=True)
                if ocr_context:
                    image_context = ocr_context
        
        # 处理需要视觉模型识别的文档
        if docs_need_vision and default_vision_model:
            
            # 如果有指定视觉模型的 provider,切换到该 provider
            if vision_provider_id:
                vision_provider = crud.get_provider(db, vision_provider_id)
                if vision_provider:
                    ai_manager.set_provider(
                        api_base=vision_provider.api_base,
                        api_key=vision_provider.api_key,
                        default_model=default_vision_model
                    )
            
            doc_context = _recognize_pdf_with_vision_model(db, docs_need_vision, default_vision_model)
            
            # 恢复原来的 provider
            if vision_provider_id and provider:
                ai_manager.set_provider(
                    api_base=provider.api_base,
                    api_key=provider.api_key,
                    default_model=model
                )
        
        # 处理需要OCR识别的文档（场景2：模型不支持视觉且未勾选眼睛按钮）
        if docs_need_ocr:
            ocr_doc_context = _recognize_docs_with_ocr(docs_need_ocr)
            if ocr_doc_context:
                if doc_context:
                    doc_context = doc_context + "\n\n" + ocr_doc_context
                else:
                    doc_context = ocr_doc_context
    
    # 合并文件上下文(非流式模式)
    all_contexts = [c for c in [file_context, image_context, doc_context] if c]
    all_file_context = "\n\n".join(all_contexts) if all_contexts else ""
    
    # 添加当前用户消息(如果有文件,将文件内容作为上下文)
    if all_file_context:
        user_content = f"{user_text}\n\n---\n以下是用户上传的文件内容,请参考:\n{all_file_context}"
    else:
        user_content = user_text
    
    # 如果当前模型支持视觉且有图片,构建多模态消息
    # 场景1：模型支持视觉时，图片直接发给AI，没有文字的文档转图片也发给AI
    if model_supports_vision and (image_files or docs_need_vision):
        import base64
        from PIL import Image
        import io
        
        content_parts = [{"type": "text", "text": user_content}]
        
        # 添加图片
        for img_info in image_files:
            try:
                with open(img_info["filepath"], "rb") as f:
                    image_data = base64.b64encode(f.read()).decode("utf-8")
                
                ext = os.path.splitext(img_info["filename"])[1].lower()
                mime_map = {
                    '.png': 'image/png',
                    '.jpg': 'image/jpeg',
                    '.jpeg': 'image/jpeg',
                    '.gif': 'image/gif',
                    '.bmp': 'image/bmp',
                    '.webp': 'image/webp'
                }
                mime_type = mime_map.get(ext, 'image/png')
                
                content_parts.append({
                    "type": "image_url",
                    "image_url": {
                        "url": f"data:{mime_type};base64,{image_data}"
                    }
                })
            except Exception:
                pass
        
        # 添加没有文字的文档（转为图片，两页合并一张）
        if docs_need_vision:
            for doc_info in docs_need_vision:
                try:
                    filepath = doc_info["filepath"]
                    filename = doc_info["filename"]
                    ext = os.path.splitext(filename)[1].lower()
                    
                    images = []
                    if ext == '.pdf':
                        from pdf2image import convert_from_path
                        import platform
                        poppler_path = None
                        if platform.system() == "Windows":
                            poppler_path = r"C:\poppler\poppler-24.08.0\Library\bin"
                            if not os.path.exists(poppler_path):
                                poppler_path = None
                        images = convert_from_path(filepath, first_page=1, last_page=10, dpi=150, poppler_path=poppler_path)
                    elif ext in ['.ppt', '.pptx']:
                        images = _convert_ppt_to_images(filepath)
                    elif ext in ['.doc', '.docx']:
                        images = _convert_word_to_images(filepath)
                    
                    if images:
                        # 两页合并一张
                        for i in range(0, len(images), 2):
                            if i + 1 < len(images):
                                img1, img2 = images[i], images[i + 1]
                                max_width = max(img1.width, img2.width)
                                if img1.width != max_width:
                                    ratio = max_width / img1.width
                                    img1 = img1.resize((max_width, int(img1.height * ratio)), Image.Resampling.LANCZOS)
                                if img2.width != max_width:
                                    ratio = max_width / img2.width
                                    img2 = img2.resize((max_width, int(img2.height * ratio)), Image.Resampling.LANCZOS)
                                total_height = img1.height + img2.height + 20
                                merged = Image.new('RGB', (max_width, total_height), 'white')
                                merged.paste(img1, (0, 0))
                                merged.paste(img2, (0, img1.height + 20))
                                img = merged
                            else:
                                img = images[i]
                            
                            buffer = io.BytesIO()
                            img.save(buffer, format='PNG', optimize=True)
                            image_data = base64.b64encode(buffer.getvalue()).decode("utf-8")
                            content_parts.append({
                                "type": "image_url",
                                "image_url": {
                                    "url": f"data:image/png;base64,{image_data}"
                                }
                            })
                except Exception as e:
                    chat_logger.warning(f"文档转图片失败 {filename}: {e}")
        
        messages.append({"role": "user", "content": content_parts})
    else:
        messages.append({"role": "user", "content": user_content})

    # 优化上下文，限制对话轮数
    messages = ContextManager.optimize_messages(messages, max_turns=6)

    # 如果启用了联网搜索，添加系统提示
    web_flag = (
        enable_web_search
        if enable_web_search is not None
        else conversation.enable_web_search
    )
    if web_flag:
        search_source = web_search_source or "duckduckgo"
        system_prompt = f"你可以使用 web_search 工具查询最新信息和实时数据。默认搜索源：{search_source}。建议：使用精准的搜索关键词，尽量一次搜索获取足够信息。"
        messages.insert(0, {"role": "system", "content": system_prompt})

    # 如果启用了 MCP 工具，添加系统提示告诉 AI 可用的工具
    mcp_flag = (
        enable_mcp
        if enable_mcp is not None
        else conversation.enable_mcp
    )
    if mcp_flag:
        mcp_tools = mcp_client.get_all_tools()
        if mcp_tools:
            tool_descriptions = []
            for tool in mcp_tools:
                func = tool.get('function', {})
                tool_name = func.get('name', '')
                tool_desc = func.get('description', '')
                tool_descriptions.append(f"- {tool_name}: {tool_desc[:150]}")
            
            mcp_system_prompt = f"""你可以使用以下工具来完成用户的请求：

{chr(10).join(tool_descriptions)}

当用户需要执行相关操作时，请积极调用工具完成任务。工具调用时使用完整的工具名称。"""
            messages.insert(0, {"role": "system", "content": mcp_system_prompt})

    # 4. 智能选择工具，减少不必要的工具定义
    conversation_tools = {
        'knowledge_base': enable_knowledge_base if enable_knowledge_base is not None else conversation.enable_knowledge_base,
        'mcp': enable_mcp if enable_mcp is not None else conversation.enable_mcp,
        'web_search': enable_web_search if enable_web_search is not None else conversation.enable_web_search
    }
    
    smart_tools = ContextManager.should_enable_tools(user_text, conversation_tools)
    
    tools_list = _build_tools_for_conversation(
        conversation,
        enable_knowledge_base=smart_tools['knowledge_base'],
        enable_mcp=smart_tools['mcp'],
        enable_web_search=smart_tools['web_search'],
    )

    # 记录聊天上下文
    logger.log_chat_context(messages, tools_list)

    # 如果没有任何工具，就走普通 chat；否则走 run_with_tools
    use_tools = bool(tools_list)

    # 5. 调用大模型
    if not stream:
        try:
            # 记录AI API调用
            logger.log_ai_api_call(
                api_base=ai_manager._provider.api_base,
                model=model or ai_manager._provider.default_model,
                messages_count=len(messages),
                tools_count=len(tools_list),
                stream=False
            )
            
            if use_tools:
                # 执行带工具的对话，包括工具调用循环
                content, token_info = _execute_chat_with_tools(
                    messages, tools_list, model, conversation_id, db
                )
            else:
                result = ai_manager.chat(messages, model=model, stream=False)
                content = result["content"]
                token_info = {
                    "model": result["model"],
                    "input_tokens": result["input_tokens"],
                    "output_tokens": result["output_tokens"],
                    "total_tokens": result["total_tokens"]
                }

            # 记录token使用情况
            logger.log_token_usage(
                model=token_info["model"],
                input_tokens=token_info["input_tokens"],
                output_tokens=token_info["output_tokens"],
                total_tokens=token_info["total_tokens"],
                estimated=token_info.get("estimated", False)
            )

            assistant_msg = crud.create_message(db, conversation_id, "assistant", content, token_info)
            logger.log_database_operation("CREATE", "messages", assistant_msg.id, {
                "role": "assistant",
                "content_length": len(content),
                "token_info": token_info
            })
            
            # 记录整体性能
            total_time = (datetime.now() - start_time).total_seconds()
            logger.log_performance("聊天完成", total_time, {
                "conversation_id": conversation_id,
                "use_tools": use_tools,
                "message_length": len(user_text),
                "response_length": len(content)
            })
            
            return {
                "user_message": {
                    "id": user_msg.id,
                    "role": user_msg.role,
                    "content": user_msg.content,
                    "created_at": user_msg.created_at.isoformat() if user_msg.created_at else None,
                },
                "user_message_id": user_msg.id,  # 明确返回用户消息ID，用于前端确认消息已保存
                "assistant_message": {
                    "id": assistant_msg.id,
                    "role": assistant_msg.role,
                    "content": assistant_msg.content,
                    "created_at": assistant_msg.created_at.isoformat() if assistant_msg.created_at else None,
                    "model": assistant_msg.model,
                    "input_tokens": assistant_msg.input_tokens,
                    "output_tokens": assistant_msg.output_tokens,
                    "total_tokens": assistant_msg.total_tokens,
                },
                "token_info": token_info,
            }
        except Exception as e:
            logger.log_error(e, "AI调用失败", {
                "conversation_id": conversation_id,
                "model": model,
                "use_tools": use_tools,
                "messages_count": len(messages),
                "tools_count": len(tools_list)
            })
            raise HTTPException(status_code=500, detail=f"AI调用失败: {str(e)}")

    # 流式:返回 StreamingResponse，最终拼接完整文本写入 DB
    def event_stream():
        nonlocal messages  # 需要修改外部的 messages 变量
        accumulated = []
        token_info = None
        vision_content_parts = []  # 收集视觉识别内容
        message_events = []  # 统一的消息事件流，按时间顺序记录所有事件
        
        # 辅助函数:添加带时间戳的事件
        import time
        def add_event(event_type: str, content):
            message_events.append({
                "type": event_type,
                "content": content,
                "timestamp": time.time()
            })
        
        # 记录流式输出开始
        chat_logger.info(f"[STREAM] 开始流式输出，对话ID: {conversation_id}, 模型: {model}")
        
        # 首先发送 ack 事件，确认用户消息已保存
        yield f"event: ack\ndata: {{\"user_message_id\": {user_msg.id}}}\n\n"
        
        # 流式模式下执行视觉识别(如果需要)
        stream_image_context = ""
        stream_doc_context = ""
        
        # 处理图片
        if image_files and not model_supports_vision:
            if force_vision_recognition and default_vision_model:
                # 用户强制启用视觉识别
                if vision_provider_id:
                    vision_provider = crud.get_provider(db, vision_provider_id)
                    if vision_provider:
                        ai_manager.set_provider(
                            api_base=vision_provider.api_base,
                            api_key=vision_provider.api_key,
                            default_model=default_vision_model
                        )
                
                image_results = []
                for event in _recognize_images_with_vision_model_stream(db, image_files, default_vision_model):
                    if event["type"] == "start":
                        yield f"event: vision_start\ndata: {json.dumps({'model': event['model'], 'total': event['total'], 'file_type': event['file_type'], 'message': '正在进行图片识别...'}, ensure_ascii=False)}\n\n"
                    elif event["type"] == "progress":
                        yield f"event: vision_progress\ndata: {json.dumps({'message': event['message']}, ensure_ascii=False)}\n\n"
                    elif event["type"] == "chunk":
                        yield f"event: vision_chunk\ndata: {json.dumps(event['content'], ensure_ascii=False)}\n\n"
                    elif event["type"] == "result":
                        image_results.append(event["content"])
                    elif event["type"] == "end":
                        yield f"event: vision_end\ndata: {json.dumps({'file_type': 'image'}, ensure_ascii=False)}\n\n"
                stream_image_context = "\n\n".join(image_results) if image_results else ""
                if stream_image_context:
                    vision_content_parts.append(stream_image_context)
                    add_event("vision", stream_image_context)
                
                if vision_provider_id and provider:
                    ai_manager.set_provider(
                        api_base=provider.api_base,
                        api_key=provider.api_key,
                        default_model=model
                    )
            else:
                # 默认使用 OCR
                # 发送开始事件
                yield f"event: vision_start\ndata: {json.dumps({'model': '本地OCR', 'total': len(image_files), 'file_type': 'image', 'message': '正在进行图片识别...'}, ensure_ascii=False)}\n\n"
                
                ocr_context, _ = _recognize_images_with_ocr(image_files, use_ocr=True)
                if ocr_context:
                    yield f"event: vision_progress\ndata: {json.dumps({'message': 'OCR识别完成'}, ensure_ascii=False)}\n\n"
                    for line in ocr_context.split('\n'):
                        if line.strip():
                            yield f"event: vision_chunk\ndata: {json.dumps(line + chr(10), ensure_ascii=False)}\n\n"
                    yield f"event: vision_end\ndata: {json.dumps({'file_type': 'image'}, ensure_ascii=False)}\n\n"
                    stream_image_context = ocr_context
                    vision_content_parts.append(ocr_context)
                    add_event("vision", ocr_context)
                else:
                    yield f"event: vision_end\ndata: {json.dumps({'file_type': 'image', 'message': '未识别到文字'}, ensure_ascii=False)}\n\n"
        
        # 处理需要视觉模型识别的文档（场景1和场景3）
        if docs_need_vision and default_vision_model and not model_supports_vision:
            # 场景3：模型不支持视觉，勾选了眼睛按钮，用视觉模型识别文档
            # 如果有指定视觉模型的 provider，切换到该 provider
            if vision_provider_id:
                vision_provider = crud.get_provider(db, vision_provider_id)
                if vision_provider:
                    ai_manager.set_provider(
                        api_base=vision_provider.api_base,
                        api_key=vision_provider.api_key,
                        default_model=default_vision_model
                    )
            
            doc_results = []
            for event in _recognize_docs_with_vision_model_stream(db, docs_need_vision, default_vision_model):
                if event["type"] == "start":
                    yield f"event: vision_start\ndata: {json.dumps({'model': event['model'], 'total': event['total'], 'file_type': event['file_type'], 'message': '正在进行文档识别...'}, ensure_ascii=False)}\n\n"
                elif event["type"] == "progress":
                    yield f"event: vision_progress\ndata: {json.dumps({'message': event['message']}, ensure_ascii=False)}\n\n"
                elif event["type"] == "chunk":
                    yield f"event: vision_chunk\ndata: {json.dumps(event['content'], ensure_ascii=False)}\n\n"
                elif event["type"] == "result":
                    doc_results.append(event["content"])
                elif event["type"] == "end":
                    yield f"event: vision_end\ndata: {json.dumps({'file_type': 'document'}, ensure_ascii=False)}\n\n"
            stream_doc_context = "\n\n".join(doc_results) if doc_results else ""
            if stream_doc_context:
                vision_content_parts.append(stream_doc_context)
                add_event("vision", stream_doc_context)
            
            # 恢复原来的 provider
            if vision_provider_id and provider:
                ai_manager.set_provider(
                    api_base=provider.api_base,
                    api_key=provider.api_key,
                    default_model=model
                )
        
        # 处理需要OCR识别的文档（场景2：模型不支持视觉且未勾选眼睛按钮）
        if docs_need_ocr:
            doc_ocr_results = []
            for event in _recognize_docs_with_ocr_stream(docs_need_ocr):
                if event["type"] == "start":
                    yield f"event: vision_start\ndata: {json.dumps({'model': event['model'], 'total': event['total'], 'file_type': event['file_type'], 'message': '正在OCR识别文档...'}, ensure_ascii=False)}\n\n"
                elif event["type"] == "progress":
                    yield f"event: vision_progress\ndata: {json.dumps({'message': event['message']}, ensure_ascii=False)}\n\n"
                elif event["type"] == "chunk":
                    yield f"event: vision_chunk\ndata: {json.dumps(event['content'], ensure_ascii=False)}\n\n"
                elif event["type"] == "result" and event.get("content"):
                    doc_ocr_results.append(event["content"])
                elif event["type"] == "end":
                    yield f"event: vision_end\ndata: {json.dumps({'file_type': 'document'}, ensure_ascii=False)}\n\n"
            doc_ocr_context = "\n\n".join(doc_ocr_results) if doc_ocr_results else ""
            if doc_ocr_context:
                if stream_doc_context:
                    stream_doc_context = stream_doc_context + "\n\n" + doc_ocr_context
                else:
                    stream_doc_context = doc_ocr_context
                vision_content_parts.append(doc_ocr_context)
                add_event("vision", doc_ocr_context)
        
        # 如果有视觉识别结果，需要更新消息内容
        if stream_image_context or stream_doc_context:
            extra_contexts = [c for c in [stream_image_context, stream_doc_context] if c]
            extra_context = "\n\n".join(extra_contexts)
            
            # 更新最后一条用户消息的内容
            if messages and messages[-1]["role"] == "user":
                last_msg = messages[-1]
                if isinstance(last_msg["content"], str):
                    if file_context:
                        # 已有文件上下文，追加视觉识别结果
                        last_msg["content"] = f"{last_msg['content']}\n\n{extra_context}"
                    else:
                        # 没有文件上下文，添加视觉识别结果
                        last_msg["content"] = f"{last_msg['content']}\n\n---\n以下是用户上传的文件内容，请参考:\n{extra_context}"
                elif isinstance(last_msg["content"], list):
                    # 多模态消息，更新第一个文本部分
                    for part in last_msg["content"]:
                        if part.get("type") == "text":
                            if file_context:
                                part["text"] = f"{part['text']}\n\n{extra_context}"
                            else:
                                part["text"] = f"{part['text']}\n\n---\n以下是用户上传的文件内容，请参考:\n{extra_context}"
                            break
        
        if use_tools:
            # 流式 + tools:灵活的工具调用和深度思考交替流程
            try:
                chat_logger.info(f"[STREAM] 使用工具模式")
                
                # 判断启用了哪些工具，发送对应提示
                kb_enabled = smart_tools.get('knowledge_base', False)
                web_enabled = smart_tools.get('web_search', False)
                mcp_enabled = smart_tools.get('mcp', False)
                
                if kb_enabled:
                    yield f"event: tool_start\ndata: {{\"status\": \"search_knowledge\", \"message\": \"正在查询知识库...\"}}\n\n"
                elif web_enabled:
                    yield f"event: tool_start\ndata: {{\"status\": \"web_search\", \"message\": \"正在联网搜索...\"}}\n\n"
                elif mcp_enabled:
                    yield f"event: tool_start\ndata: {{\"status\": \"mcp\", \"message\": \"正在调用工具...\"}}\n\n"
                else:
                    yield f"event: tool_start\ndata: {{\"status\": \"thinking\", \"message\": \"正在分析问题...\"}}\n\n"
                
                current_messages = messages.copy()
                tool_calls_info = []
                thinking_content = []
                total_input_tokens = 0
                total_output_tokens = 0
                max_iterations = 3  # 限制工具调用次数，避免过多消耗
                
                # 第一阶段:非深度思考模式下的工具调用循环
                for iteration in range(max_iterations):
                    # 调用模型(非流式，不启用深度思考)
                    # 打印消息结构用于调试
                    for i, msg in enumerate(current_messages):
                        role = msg.get('role', 'unknown')
                        content = msg.get('content', '')
                        tool_calls = msg.get('tool_calls')
                        tool_call_id = msg.get('tool_call_id')
                        content_preview = str(content)[:50] if content else 'None'
                    
                    data = ai_manager.run_with_tools(current_messages, tools=tools_list, model=model, stream=False)
                    
                    usage = data.get("usage", {})
                    total_input_tokens += usage.get("prompt_tokens", 0)
                    total_output_tokens += usage.get("completion_tokens", 0)
                    
                    message = data["choices"][0]["message"]
                    
                    tool_calls = message.get("tool_calls")
                    if not tool_calls:
                        # 没有工具调用，发送提示并进入深度思考阶段
                        yield f"event: tool_end\ndata: {{\"status\": \"skipped\", \"message\": \"模型未调用工具\", \"tools\": []}}\n\n"
                        break
                    
                    # 有工具调用，把消息加入历史
                    # 注意：需要确保 message 格式正确，某些 API 可能返回额外字段
                    assistant_msg = {
                        "role": "assistant",
                        "content": message.get("content") or "",
                        "tool_calls": message.get("tool_calls")
                    }
                    current_messages.append(assistant_msg)
                    
                    # 执行工具调用
                    for tool_call in tool_calls:
                        function_name = tool_call["function"]["name"]
                        function_args = json.loads(tool_call["function"]["arguments"])
                        
                        # 发送工具调用进度 - 开始
                        # 处理 MCP 工具名称显示
                        if function_name.startswith("mcp_"):
                            parts = function_name.split("_", 2)
                            if len(parts) >= 3:
                                tool_display_name = f"MCP:{parts[1]}:{parts[2]}"
                            else:
                                tool_display_name = function_name
                        else:
                            tool_display_name = {
                                "search_knowledge": "知识库搜索",
                                "web_search": "联网搜索",
                                "get_local_time": "获取时间",
                                "calculate_expression": "计算器"
                            }.get(function_name, function_name)
                        
                        # 构建搜索参数显示
                        if function_name == "search_knowledge":
                            query = function_args.get("query", "")
                            top_k = function_args.get("top_k", 5)
                            yield f"event: tool_progress\ndata: {{\"tool\": \"{function_name}\", \"stage\": \"start\", \"message\": \"正在搜索: {query}\"}}\n\n"
                        elif function_name == "web_search":
                            query = function_args.get("query", "")
                            yield f"event: tool_progress\ndata: {{\"tool\": \"{function_name}\", \"stage\": \"start\", \"message\": \"正在搜索: {query}\"}}\n\n"
                        elif function_name.startswith("mcp_"):
                            # MCP 工具调用
                            yield f"event: tool_progress\ndata: {{\"tool\": \"{function_name}\", \"stage\": \"start\", \"message\": \"正在调用 {tool_display_name}...\"}}\n\n"
                        else:
                            yield f"event: tool_progress\ndata: {{\"tool\": \"{function_name}\", \"stage\": \"start\", \"message\": \"正在执行 {tool_display_name}...\"}}\n\n"
                        
                        tool_info = {"name": function_name, "args": function_args, "status": "running"}
                        tool_calls_info.append(tool_info)
                        
                        try:
                            result = _execute_tool(function_name, function_args, conversation_id, db)
                            tool_info["status"] = "success"
                            # 提取结果预览
                            result_preview = result[:150] + "..." if len(result) > 150 else result
                            tool_info["result_preview"] = result_preview
                            
                            # 发送工具调用进度 - 完成
                            yield f"event: tool_progress\ndata: {{\"tool\": \"{function_name}\", \"stage\": \"done\", \"message\": \"✓ 搜索完成\", \"preview\": {json.dumps(result_preview, ensure_ascii=False)}}}\n\n"
                            
                            # 记录工具调用事件
                            add_event("tool_call", tool_info.copy())
                        except Exception as e:
                            result = f"工具执行失败: {str(e)}"
                            tool_info["status"] = "error"
                            tool_info["error"] = str(e)
                            yield f"event: tool_progress\ndata: {{\"tool\": \"{function_name}\", \"stage\": \"error\", \"message\": \"✗ 执行失败: {str(e)}\"}}\n\n"
                            
                            # 记录失败的工具调用事件
                            add_event("tool_call", tool_info.copy())
                        
                        current_messages.append({
                            "role": "tool",
                            "tool_call_id": tool_call["id"],
                            "content": result
                        })
                
                # 发送工具调用完成提示
                if tool_calls_info:
                    tool_end_data = f"event: tool_end\ndata: {{\"status\": \"done\", \"tools\": {json.dumps(tool_calls_info, ensure_ascii=False)}}}\n\n"
                    yield tool_end_data
                    # 重置工具调用信息，避免累计到第二阶段
                    tool_calls_info = []
                
                # 第二阶段:深度思考模式下的流式生成(支持模型自主决定是否继续调用工具)
                
                is_thinking_done = False
                final_response_iterations = 0
                max_final_iterations = 5  # 深度思考阶段最多允许的额外工具调用轮数
                
                # 初始化 XML 工具调用相关变量(在循环外)
                xml_tool_buffer = ""
                in_xml_tool_call = False
                has_xml_tool_call = False
                
                while final_response_iterations < max_final_iterations:
                    final_response_iterations += 1
                    
                    # 如果启用深度思考，发送思考开始提示
                    if enable_thinking and not is_thinking_done:
                        yield f"event: thinking_start\ndata: {{\"status\": \"thinking\", \"message\": \"正在深度思考...\"}}\n\n"
                    
                    # 重置待输出内容(每次迭代都重置)
                    pending_output = []  # 待输出的内容，用于延迟输出以检测 XML 工具调用
                    
                    # 标记是否已经有正文内容(用于判断 reasoning_content 是否应该作为正文)
                    has_real_content = False
                    thinking_buffer = []  # 用于累积思考内容，检测XML工具调用
                    
                    for chunk in ai_manager.chat(current_messages, model=model, stream=True, enable_thinking=enable_thinking):
                        if isinstance(chunk, dict):
                            chunk_type = chunk.get("type", "")
                            
                            if chunk_type == "usage":
                                usage = chunk.get("usage", {})
                                total_input_tokens += usage.get("prompt_tokens", 0)
                                total_output_tokens += usage.get("completion_tokens", 0)
                                continue
                            
                            # 处理思考内容
                            if chunk_type == "thinking":
                                thinking = chunk.get("content", "")
                                if thinking:
                                    thinking_content.append(thinking)
                                    thinking_buffer.append(thinking)
                                    
                                    # 检测思考内容中是否有XML工具调用
                                    thinking_so_far = "".join(thinking_buffer)
                                    import re
                                    # 支持多种格式: <function_calls>, <| DSML | function_calls>, <function_calls> 等
                                    fc_match = re.search(r'<[\s\|]*(?:DSML\s*\|)?\s*(?:antml:)?function_calls\s*>', thinking_so_far, re.IGNORECASE)
                                    if fc_match:
                                        # 发送思考结束事件(只发送 <function_calls> 之前的内容)
                                        fc_start = fc_match.start()
                                        thinking_before_fc = thinking_so_far[:fc_start]
                                        if thinking_before_fc.strip():
                                            yield f"event: thinking_end\ndata: {{\"thinking\": {json.dumps(thinking_before_fc, ensure_ascii=False)}}}\n\n"
                                            add_event("thinking", thinking_before_fc)
                                        is_thinking_done = True
                                        
                                        # 开始收集XML工具调用
                                        in_xml_tool_call = True
                                        xml_tool_buffer = thinking_so_far[fc_start:]
                                        thinking_buffer = []  # 清空缓冲区
                                        
                                        # 检测是否已经结束
                                        # 支持多种格式: </function_calls>, </| DSML | function_calls>
                                        if re.search(r'</[\s\|]*(?:DSML\s*\|)?\s*(?:antml:)?function_calls\s*>', xml_tool_buffer, re.IGNORECASE):
                                            in_xml_tool_call = False
                                            has_xml_tool_call = True
                                        continue
                                    
                                    # 如果正在收集XML工具调用，继续收集
                                    if in_xml_tool_call:
                                        xml_tool_buffer += thinking
                                        # 支持多种格式
                                        if re.search(r'</[\s\|]*(?:DSML\s*\|)?\s*(?:antml:)?function_calls\s*>', xml_tool_buffer, re.IGNORECASE):
                                            in_xml_tool_call = False
                                            has_xml_tool_call = True
                                        continue
                                    
                                    # 检测是否可能是 XML 工具调用的开始（需要等待更多内容）
                                    # 检查最后是否有未完成的 < 标签
                                    potential_xml_in_thinking = False
                                    if '<' in thinking_so_far:
                                        last_lt_pos = thinking_so_far.rfind('<')
                                        remaining = thinking_so_far[last_lt_pos:].lower().replace(' ', '').replace('\n', '')
                                        # 检查是否可能是 <function_calls> 或 <| DSML | function_calls> 的开始
                                        possible_starts = ['<function_calls>', '<|dsml|function_calls>', '<|', '<f', '<fu', '<fun', '<func', '<funct', '<functi', '<functio', '<function', '<function_', '<function_c', '<function_ca', '<function_cal', '<function_call', '<function_calls', '<ds', '<dsm', '<dsml']
                                        for ps in possible_starts:
                                            if ps.startswith(remaining) or remaining.startswith(ps.rstrip('>')):
                                                potential_xml_in_thinking = True
                                                break
                                    
                                    if potential_xml_in_thinking:
                                        # 可能是 XML 开始，暂不发送，继续缓冲
                                        continue
                                    
                                    # 不是 XML，发送思考内容
                                    # 但只发送到最后一个 < 之前的内容（如果有的话）
                                    if '<' in thinking_so_far:
                                        last_lt_pos = thinking_so_far.rfind('<')
                                        safe_content = thinking_so_far[:last_lt_pos]
                                        if safe_content.strip():
                                            yield f"event: thinking\ndata: {json.dumps(safe_content, ensure_ascii=False)}\n\n"
                                        # 保留 < 之后的内容继续缓冲
                                        thinking_buffer = [thinking_so_far[last_lt_pos:]]
                                    else:
                                        yield f"event: thinking\ndata: {json.dumps(thinking_so_far, ensure_ascii=False)}\n\n"
                                        thinking_buffer = []  # 清空缓冲区
                                continue
                            
                            # 处理正文内容
                            if chunk_type == "content":
                                delta = chunk.get("content", "")
                                if delta:
                                    has_real_content = True
                            else:
                                delta = chunk.get("content", "")
                        else:
                            delta = str(chunk)
                        
                        if delta:
                            # 如果已经在收集 XML 工具调用，继续收集
                            if in_xml_tool_call:
                                xml_tool_buffer += delta
                                # 检测工具调用结束(支持多种格式)
                                import re
                                if re.search(r'</[\s\|]*(?:DSML\s*\|)?\s*(?:antml:)?function_calls\s*>', xml_tool_buffer, re.IGNORECASE):
                                    in_xml_tool_call = False
                                    has_xml_tool_call = True
                                continue
                            
                            # 累积内容用于检测 XML 工具调用开始
                            pending_output.append(delta)
                            pending_content = "".join(pending_output)
                            
                            # 检测是否可能是 XML 工具调用的开始
                            # 检查是否包含 < 且可能是 <function_calls> 的开始
                            potential_xml_start = False
                            if "<" in pending_content:
                                # 检查是否是 <function_calls> 或其变体的部分匹配
                                last_lt_pos = pending_content.rfind("<")
                                remaining = pending_content[last_lt_pos:].lower().replace(" ", "").replace("\n", "")
                                # 检查是否是各种变体的开始（包括 <function_calls>）
                                # 完整的目标标签列表
                                target_tags = [
                                    '<function_calls>',
                                    '<function_calls>',
                                    '<|dsml|function_calls>',
                                ]
                                # 生成所有可能的前缀
                                possible_prefixes = set()
                                for tag in target_tags:
                                    for i in range(1, len(tag)):
                                        possible_prefixes.add(tag[:i].lower())
                                
                                # 检查 remaining 是否是某个目标标签的前缀
                                for tag in target_tags:
                                    tag_lower = tag.lower()
                                    if tag_lower.startswith(remaining) and len(remaining) < len(tag_lower):
                                        potential_xml_start = True
                                        break
                                
                                # 如果缓冲内容太长（超过50字符）还没匹配到，说明不是 XML
                                if potential_xml_start and len(remaining) > 50:
                                    potential_xml_start = False
                            
                            # 检测完整的 <function_calls> 标签(支持多种格式)
                            # 使用正则表达式进行更灵活的匹配
                            import re
                            fc_match = re.search(r'<[\s\|]*(?:DSML\s*\|)?\s*(?:antml:)?function_calls\s*>', pending_content, re.IGNORECASE)
                            if fc_match:
                                in_xml_tool_call = True
                                fc_start = fc_match.start()
                                # 输出 <function_calls> 之前的内容
                                if fc_start > 0:
                                    before_fc = pending_content[:fc_start]
                                    if before_fc.strip():
                                        # 发送思考结束事件
                                        if enable_thinking and not is_thinking_done:
                                            is_thinking_done = True
                                            full_thinking = "".join(thinking_content) if thinking_content else ""
                                            yield f"event: thinking_end\ndata: {{\"thinking\": {json.dumps(full_thinking, ensure_ascii=False)}}}\n\n"
                                            # 记录思考事件
                                            if full_thinking:
                                                add_event("thinking", full_thinking)
                                        yield f"data: {json.dumps(before_fc, ensure_ascii=False)}\n\n"
                                        accumulated.append(before_fc)
                                        # 记录正文事件
                                        add_event("text", before_fc)
                                # 开始收集工具调用内容
                                xml_tool_buffer = pending_content[fc_start:]
                                pending_output = []
                                # 检测工具调用是否已经结束(支持多种格式)
                                fc_end_match = re.search(r'</[\s\|]*(?:DSML\s*\|)?\s*(?:antml:)?function_calls\s*>', xml_tool_buffer, re.IGNORECASE)
                                if fc_end_match:
                                    in_xml_tool_call = False
                                    has_xml_tool_call = True
                                continue
                            
                            # 如果可能是 XML 开始，继续等待更多内容
                            if potential_xml_start:
                                continue
                            
                            # 不是 XML 工具调用，正常输出
                            # 如果启用了深度思考且还没发送结束事件，在开始输出正文时发送
                            if enable_thinking and not is_thinking_done:
                                is_thinking_done = True
                                full_thinking = "".join(thinking_content) if thinking_content else ""
                                yield f"event: thinking_end\ndata: {{\"thinking\": {json.dumps(full_thinking, ensure_ascii=False)}}}\n\n"
                                # 记录思考事件
                                if full_thinking:
                                    add_event("thinking", full_thinking)
                            
                            # 输出所有待输出的内容
                            output_content = "".join(pending_output)
                            if output_content:
                                accumulated.append(output_content)
                                yield f"data: {json.dumps(output_content, ensure_ascii=False)}\n\n"
                                # 记录正文事件
                                add_event("text", output_content)
                            pending_output = []
                    
                    # 流结束后，检查思考内容中是否有未处理的 XML 工具调用
                    if thinking_buffer and not in_xml_tool_call and not has_xml_tool_call:
                        thinking_so_far = "".join(thinking_buffer)
                        import re
                        # 支持多种格式
                        fc_match = re.search(r'<[\s\|]*(?:DSML\s*\|)?\s*(?:antml:)?function_calls\s*>', thinking_so_far, re.IGNORECASE)
                        if fc_match:
                            fc_start = fc_match.start()
                            xml_tool_buffer = thinking_so_far[fc_start:]
                            if re.search(r'</[\s\|]*(?:DSML\s*\|)?\s*(?:antml:)?function_calls\s*>', xml_tool_buffer, re.IGNORECASE):
                                has_xml_tool_call = True
                                # 发送思考结束事件(只发送 <function_calls> 之前的内容)
                                thinking_before_fc = thinking_so_far[:fc_start]
                                if thinking_before_fc.strip() and not is_thinking_done:
                                    yield f"event: thinking_end\ndata: {{\"thinking\": {json.dumps(thinking_before_fc, ensure_ascii=False)}}}\n\n"
                                    add_event("thinking", thinking_before_fc)
                                    is_thinking_done = True
                    
                    # 流结束后，处理剩余内容
                    # 如果正在收集 XML 工具调用，检查是否完整
                    if in_xml_tool_call:
                        import re
                        if re.search(r'</[\s\|]*(?:DSML\s*\|)?\s*(?:antml:)?function_calls\s*>', xml_tool_buffer, re.IGNORECASE):
                            in_xml_tool_call = False
                            has_xml_tool_call = True
                        else:
                            # XML 不完整，作为普通内容输出
                            if enable_thinking and not is_thinking_done:
                                is_thinking_done = True
                                full_thinking = "".join(thinking_content) if thinking_content else ""
                                yield f"event: thinking_end\ndata: {{\"thinking\": {json.dumps(full_thinking, ensure_ascii=False)}}}\n\n"
                                if full_thinking:
                                    add_event("thinking", full_thinking)
                            accumulated.append(xml_tool_buffer)
                            yield f"data: {json.dumps(xml_tool_buffer, ensure_ascii=False)}\n\n"
                            add_event("text", xml_tool_buffer)
                            in_xml_tool_call = False
                    
                    # 输出剩余的待输出内容
                    if pending_output and not in_xml_tool_call:
                        if enable_thinking and not is_thinking_done:
                            is_thinking_done = True
                            full_thinking = "".join(thinking_content) if thinking_content else ""
                            yield f"event: thinking_end\ndata: {{\"thinking\": {json.dumps(full_thinking, ensure_ascii=False)}}}\n\n"
                            # 记录思考事件
                            if full_thinking:
                                add_event("thinking", full_thinking)
                        output_content = "".join(pending_output)
                        if output_content:
                            accumulated.append(output_content)
                            yield f"data: {json.dumps(output_content, ensure_ascii=False)}\n\n"
                            # 记录正文事件
                            add_event("text", output_content)
                    
                    # 检查是否有 XML 工具调用需要执行
                    if has_xml_tool_call:
                        
                        # 解析 XML 工具调用(支持多种格式)
                        # 支持: <invoke name="...">, <| DSML | invoke name="...">, <invoke name="...">
                        invoke_pattern = r'<[\s\|]*(?:DSML\s*\|)?\s*(?:antml:)?invoke\s+name\s*=\s*["\']([^"\']+)["\']\s*>(.*?)</[\s\|]*(?:DSML\s*\|)?\s*(?:antml:)?invoke\s*>'
                        # 支持: <parameter name="...">, <| DSML | parameter name="...">, <parameter name="...">
                        param_pattern = r'<[\s\|]*(?:DSML\s*\|)?\s*(?:antml:)?parameter\s+name\s*=\s*["\']([^"\']+)["\'][^>]*>(.*?)</[\s\|]*(?:DSML\s*\|)?\s*(?:antml:)?parameter\s*>'
                        
                        xml_tool_results = []
                        xml_tool_count = len(re.findall(invoke_pattern, xml_tool_buffer, re.DOTALL | re.IGNORECASE))
                        
                        for invoke_match in re.finditer(invoke_pattern, xml_tool_buffer, re.DOTALL | re.IGNORECASE):
                            tool_name = invoke_match.group(1)
                            params_str = invoke_match.group(2)
                            
                            # 解析参数
                            params = {}
                            for param_match in re.finditer(param_pattern, params_str, re.DOTALL | re.IGNORECASE):
                                param_name = param_match.group(1)
                                param_value = param_match.group(2).strip()
                                try:
                                    params[param_name] = int(param_value)
                                except ValueError:
                                    params[param_name] = param_value
                            
                            # 发送工具调用开始提示
                            tool_messages = {
                                "search_knowledge": "正在查询知识库...",
                                "web_search": "正在联网搜索...",
                            }
                            tool_msg = tool_messages.get(tool_name, f"正在执行 {tool_name}...")
                            yield f"event: tool_start\ndata: {{\"status\": \"{tool_name}\", \"message\": \"{tool_msg}\"}}\n\n"
                            
                            # 发送工具进度开始事件
                            if tool_name == "search_knowledge":
                                query = params.get("query", "")
                                yield f"event: tool_progress\ndata: {{\"tool\": \"{tool_name}\", \"stage\": \"start\", \"message\": \"正在搜索: {query}\"}}\n\n"
                            else:
                                yield f"event: tool_progress\ndata: {{\"tool\": \"{tool_name}\", \"stage\": \"start\", \"message\": \"正在执行 {tool_name}...\"}}\n\n"
                            
                            tool_info = {"name": tool_name, "args": params, "status": "running"}
                            tool_calls_info.append(tool_info)
                            
                            try:
                                result = _execute_tool(tool_name, params, conversation_id, db)
                                tool_info["status"] = "success"
                                tool_info["result_preview"] = result[:100] + "..." if len(result) > 100 else result
                                xml_tool_results.append(f"工具 {tool_name} 执行结果:\n{result}")
                                
                                # 发送工具进度完成事件
                                result_preview = result[:50] + "..." if len(result) > 50 else result
                                yield f"event: tool_progress\ndata: {{\"tool\": \"{tool_name}\", \"stage\": \"done\", \"message\": \"✓ 执行完成\", \"preview\": {json.dumps(result_preview, ensure_ascii=False)}}}\n\n"
                                
                                # 记录工具调用事件
                                add_event("tool_call", tool_info.copy())
                            except Exception as e:
                                result = f"工具执行失败: {str(e)}"
                                tool_info["status"] = "error"
                                tool_info["error"] = str(e)
                                xml_tool_results.append(f"工具 {tool_name} 执行失败: {str(e)}")
                                
                                # 发送工具进度错误事件
                                yield f"event: tool_progress\ndata: {{\"tool\": \"{tool_name}\", \"stage\": \"error\", \"message\": \"✗ 执行失败: {str(e)}\"}}\n\n"
                                
                                # 记录失败的工具调用事件
                                add_event("tool_call", tool_info.copy())
                        
                        # 发送工具调用完成提示(只包含本轮的工具调用信息)
                        yield f"event: tool_end\ndata: {{\"status\": \"done\", \"tools\": {json.dumps(tool_calls_info, ensure_ascii=False)}}}\n\n"
                        
                        # 重置工具调用信息，避免累计
                        tool_calls_info = []
                        
                        # 将工具结果添加到消息中，继续对话
                        # 注意：不要将 XML 工具调用内容作为 assistant 消息，这可能导致 API 400 错误
                        # 而是将工具执行结果作为 system 消息添加，让模型基于结果继续回复
                        current_messages.append({
                            "role": "system",
                            "content": f"工具执行结果:\n\n" + "\n\n".join(xml_tool_results)
                        })
                        
                        # 重置状态，继续下一轮
                        is_thinking_done = False
                        thinking_content = []  # 清空思考内容，准备新一轮
                        xml_tool_buffer = ""  # 重置 XML 工具调用缓冲区
                        in_xml_tool_call = False
                        has_xml_tool_call = False
                        continue
                    else:
                        # 没有工具调用，结束循环
                        break
                
                # 如果启用了深度思考但流结束时还没发送结束事件
                if enable_thinking and not is_thinking_done:
                    full_thinking = "".join(thinking_content) if thinking_content else ""
                    yield f"event: thinking_end\ndata: {{\"thinking\": {json.dumps(full_thinking, ensure_ascii=False)}}}\n\n"
                    # 记录思考事件
                    if full_thinking:
                        add_event("thinking", full_thinking)
                
                # 特殊处理:如果没有正文内容但有思考内容
                # 检查思考内容中是否有 XML 工具调用
                if not accumulated and thinking_content:
                    full_thinking_as_content = "".join(thinking_content)
                    
                    # 检查是否有 XML 工具调用(支持多种格式)
                    import re
                    fc_match = re.search(r'<[\s\|]*(?:DSML\s*\|)?\s*(?:antml:)?function_calls\s*>', full_thinking_as_content, re.IGNORECASE)
                    fc_end_match = re.search(r'</[\s\|]*(?:DSML\s*\|)?\s*(?:antml:)?function_calls\s*>', full_thinking_as_content, re.IGNORECASE)
                    
                    if fc_match and fc_end_match:
                        has_xml_tool_call = True
                        xml_tool_buffer = full_thinking_as_content[fc_match.start():]
                        # 不输出思考内容，让工具调用逻辑处理
                    else:
                        # 没有工具调用，把思考内容作为正文输出
                        # 发送正文内容
                        yield f"data: {json.dumps(full_thinking_as_content, ensure_ascii=False)}\n\n"
                        accumulated.append(full_thinking_as_content)
                        # 记录正文事件(思考内容作为正文)
                        add_event("text", full_thinking_as_content)
                        # 清空思考内容，因为已经作为正文输出了
                        thinking_content = []
                
                token_info = {
                    "model": model or "default",
                    "input_tokens": total_input_tokens,
                    "output_tokens": total_output_tokens,
                    "total_tokens": total_input_tokens + total_output_tokens
                }
                
                chat_logger.info(f"[STREAM] 流式输出完成")
                
                # 发送token信息
                yield f"event: meta\ndata: {json.dumps(token_info)}\n\n"
                yield "data: [DONE]\n\n"
                
                # 写入数据库
                full_text = "".join(accumulated)
                try:
                    logger.log_token_usage(
                        model=token_info.get("model", model or "default"),
                        input_tokens=token_info.get("input_tokens", 0),
                        output_tokens=token_info.get("output_tokens", 0),
                        total_tokens=token_info.get("total_tokens", 0),
                        estimated=token_info.get("estimated", False),
                    )
                except Exception:
                    pass
                
                # 保存工具调用、深度思考内容、视觉识别内容和消息事件流
                tool_calls_json = json.dumps(tool_calls_info, ensure_ascii=False) if tool_calls_info else None
                full_thinking = "".join(thinking_content) if thinking_content else None
                full_vision = "\n\n".join(vision_content_parts) if vision_content_parts else None
                message_events_json = json.dumps(message_events, ensure_ascii=False) if message_events else None
                crud.create_message(db, conversation_id, "assistant", full_text, token_info, 
                                   tool_calls=tool_calls_json, thinking_content=full_thinking,
                                   vision_content=full_vision, message_events=message_events_json)
                
            except Exception as e:
                chat_logger.error(f"[STREAM] 工具模式错误: {str(e)}")
                import traceback
                traceback.print_exc()
                yield f"data: [错误] {str(e)}\n\n"
                yield "data: [DONE]\n\n"

        else:
            try:
                chat_logger.info(f"[STREAM] 普通流式模式，深度思考: {enable_thinking}")
                chunk_count = 0
                thinking_content = []  # 存储思考内容
                is_thinking = False
                
                # 如果启用深度思考，先发送思考开始提示
                if enable_thinking:
                    yield f"event: thinking_start\ndata: {{\"status\": \"thinking\", \"message\": \"正在深度思考...\"}}\n\n"
                
                # 普通流式对话，直接消费 include_usage 终结器
                for chunk in ai_manager.chat(messages, model=model, stream=True, enable_thinking=enable_thinking):
                    if isinstance(chunk, dict):
                        if chunk.get("type") == "usage":
                            token_info = chunk.get("usage")
                            continue
                        
                        # 处理思考内容
                        if chunk.get("type") == "thinking":
                            thinking = chunk.get("content", "")
                            if thinking:
                                thinking_content.append(thinking)
                                # 发送思考内容(前端可以选择显示或隐藏)
                                yield f"event: thinking\ndata: {json.dumps(thinking, ensure_ascii=False)}\n\n"
                            continue
                        
                        delta = chunk.get("content", "")
                    else:
                        delta = str(chunk)

                    if delta:
                        # 如果之前在思考，现在开始输出正文，发送思考结束事件
                        if enable_thinking and thinking_content and not is_thinking:
                            is_thinking = True
                            full_thinking = "".join(thinking_content)
                            yield f"event: thinking_end\ndata: {{\"thinking\": {json.dumps(full_thinking, ensure_ascii=False)}}}\n\n"
                            # 记录思考事件
                            if full_thinking:
                                add_event("thinking", full_thinking)
                        
                        accumulated.append(delta)
                        chunk_count += 1
                        # 使用 JSON 编码以保留换行符(SSE 中换行符会破坏格式)
                        yield f"data: {json.dumps(delta, ensure_ascii=False)}\n\n"
                
                full_text = "".join(accumulated)
                # 记录最终正文事件(普通模式下正文是连续的)
                if full_text:
                    add_event("text", full_text)
                chat_logger.info(f"[STREAM] 流式输出完成，共 {chunk_count} 个块，总长度: {len(full_text)}")

                # 如果流式没给 usage，则估算
                if not token_info:
                    estimated_input = max(1, len(json.dumps(messages, ensure_ascii=False)) // 4)
                    estimated_output = max(1, len(full_text) // 4)
                    token_info = {
                        "model": model or "default",
                        "input_tokens": estimated_input,
                        "output_tokens": estimated_output,
                        "total_tokens": estimated_input + estimated_output,
                        "estimated": True,
                    }

                # 发送token信息
                yield f"event: meta\ndata: {json.dumps(token_info)}\n\n"
                yield "data: [DONE]\n\n"
                
                chat_logger.info(f"[STREAM] 发送 [DONE] 标记")
                
                # 完成后将完整回复写入数据库
                try:
                    logger.log_token_usage(
                        model=token_info.get("model", model or "default"),
                        input_tokens=token_info.get("input_tokens", 0),
                        output_tokens=token_info.get("output_tokens", 0),
                        total_tokens=token_info.get("total_tokens", 0),
                        estimated=token_info.get("estimated", False),
                    )
                except Exception:
                    pass

                # 保存深度思考内容、视觉识别内容和消息事件流(普通模式没有工具调用)
                full_thinking = "".join(thinking_content) if thinking_content else None
                full_vision = "\n\n".join(vision_content_parts) if vision_content_parts else None
                message_events_json = json.dumps(message_events, ensure_ascii=False) if message_events else None
                crud.create_message(db, conversation_id, "assistant", full_text, token_info,
                                   tool_calls=None, thinking_content=full_thinking,
                                   vision_content=full_vision, message_events=message_events_json)
                
            except Exception as e:
                chat_logger.error(f"[STREAM] 普通模式错误: {str(e)}")
                yield f"data: [错误] {str(e)}\n\n"
                yield "data: [DONE]\n\n"

    return StreamingResponse(event_stream(), media_type="text/event-stream")

# ========== 文件上传(对话级) ==========

UPLOAD_DIR = "uploads"
os.makedirs(UPLOAD_DIR, exist_ok=True)

@app.post("/upload")
def upload_file(
    conversation_id: int = Form(...),
    file: UploadFile = File(...),
    db: Session = Depends(get_db),
):
    conversation = crud.get_conversation(db, conversation_id)
    if not conversation:
        raise HTTPException(status_code=404, detail="Conversation not found")

    save_dir = os.path.join(UPLOAD_DIR, str(conversation_id))
    os.makedirs(save_dir, exist_ok=True)

    save_path = os.path.join(save_dir, file.filename)
    with open(save_path, "wb") as f:
        shutil.copyfileobj(file.file, f)

    record = crud.create_uploaded_file(db, conversation_id, file.filename, save_path)
    return record.to_dict()

@app.get("/conversations/{conversation_id}/files")
def list_conversation_files(
    conversation_id: int,
    db: Session = Depends(get_db),
):
    files = crud.get_uploaded_files(db, conversation_id)
    return [file.to_dict() for file in files]

@app.delete("/files/{file_id}")
def delete_conversation_file(file_id: int, db: Session = Depends(get_db)):
    file_record = crud.get_uploaded_file(db, file_id)
    if not file_record:
        raise HTTPException(status_code=404, detail="File not found")

    # 删除本地文件
    try:
        if os.path.exists(file_record.filepath):
            os.remove(file_record.filepath)
    except Exception:
        pass

    crud.delete_uploaded_file(db, file_id)
    return {"success": True}

@app.get("/files/{path:path}")
def get_file(path: str):
    file_path = os.path.join(UPLOAD_DIR, path)
    if not os.path.exists(file_path):
        raise HTTPException(status_code=404, detail="File not found")
    return FileResponse(file_path)

# ========== Provider 管理接口(新增) ==========

@app.get("/providers")
def list_providers(db: Session = Depends(get_db)):
    providers = crud.list_providers(db)
    return [provider.to_dict() for provider in providers]

@app.get("/providers/{provider_id}")
def get_provider_detail(provider_id: int, db: Session = Depends(get_db)):
    """获取单个Provider的详细信息(用于调试)"""
    provider = crud.get_provider(db, provider_id)
    if not provider:
        raise HTTPException(status_code=404, detail="Provider not found")
    # 返回是否有API Key的状态
    return {
        "id": provider.id,
        "name": provider.name,
        "api_base": provider.api_base,
        "has_api_key": bool(provider.api_key),
        "api_key_length": len(provider.api_key) if provider.api_key else 0,
        "default_model": provider.default_model,
        "models": provider.models,
        "is_default": provider.is_default,
    }

@app.post("/providers")
def create_provider(
    name: str = Form(...),
    api_base: str = Form(...),
    api_key: str = Form(...),
    default_model: str = Form(...),
    models_str: Optional[str] = Form(None),
    models_config: Optional[str] = Form(None),
    is_default: bool = Form(False),
    db: Session = Depends(get_db),
):
    try:
        # 检查是否已存在同名Provider
        existing = crud.get_provider_by_name(db, name)
        if existing:
            raise HTTPException(status_code=400, detail=f"Provider名称 '{name}' 已存在")
        
        provider = crud.create_provider(
            db,
            name=name,
            api_base=api_base,
            api_key=api_key,
            default_model=default_model,
            models_str=models_str,
            models_config=models_config,
            is_default=is_default,
        )
        return provider.to_dict()
    except HTTPException:
        raise
    except Exception as e:
        chat_logger.error(f"创建Provider失败: {e}")
        raise HTTPException(status_code=500, detail=f"创建Provider失败: {str(e)}")

@app.post("/providers/{provider_id}")
def update_provider(
    provider_id: int,
    name: Optional[str] = Form(None),
    api_base: Optional[str] = Form(None),
    api_key: Optional[str] = Form(None),
    default_model: Optional[str] = Form(None),
    models_str: Optional[str] = Form(None),
    models_config: Optional[str] = Form(None),
    is_default: Optional[bool] = Form(None),
    db: Session = Depends(get_db),
):
    # 如果api_key为空字符串，则不更新(保持原值)
    actual_api_key = api_key if api_key and api_key.strip() else None
    
    provider = crud.update_provider(
        db,
        provider_id,
        name=name,
        api_base=api_base,
        api_key=actual_api_key,
        default_model=default_model,
        models_str=models_str,
        models_config=models_config,
        is_default=is_default,
    )
    if not provider:
        raise HTTPException(status_code=404, detail="Provider not found")
    return provider.to_dict()

@app.delete("/providers/{provider_id}")
def delete_provider(provider_id: int, db: Session = Depends(get_db)):
    crud.delete_provider(db, provider_id)
    return {"success": True}

@app.get("/providers/{provider_id}/models")
def get_provider_models(provider_id: int, db: Session = Depends(get_db)):
    """获取指定Provider的模型列表"""
    provider = crud.get_provider(db, provider_id)
    if not provider:
        raise HTTPException(status_code=404, detail="Provider not found")
    
    models = []
    if provider.models:
        models = [m.strip() for m in provider.models.split(",") if m.strip()]
    else:
        models = [provider.default_model]
    
    return {
        "provider_id": provider_id,
        "provider_name": provider.name,
        "default_model": provider.default_model,
        "models": models,
    }

@app.get("/models/all")
def get_all_models(db: Session = Depends(get_db)):
    """获取所有Provider的模型列表，用于前端统一显示"""
    import json
    providers = crud.list_providers(db)
    all_models = set()
    models_caps = {}  # 存储每个模型的功能信息
    models_names = {}  # 存储每个模型的自定义显示名称
    
    # 添加全局默认模型
    all_models.update(settings.ai_models)
    
    # 添加所有Provider的模型
    for provider in providers:
        # 解析模型配置
        config = {}
        if provider.models_config:
            try:
                config = json.loads(provider.models_config)
            except:
                pass
        
        # 始终添加默认模型
        all_models.add(provider.default_model)
        
        # 添加其他模型
        if provider.models:
            provider_models = [m.strip() for m in provider.models.split(",") if m.strip()]
            all_models.update(provider_models)
            # 合并功能信息和自定义名称
            for model in provider_models:
                if model in config:
                    model_config = config[model]
                    models_caps[model] = model_config
                    # 提取自定义名称
                    if model_config.get("custom_name"):
                        models_names[model] = model_config["custom_name"]
        
        # 默认模型的功能信息和自定义名称
        if provider.default_model in config:
            model_config = config[provider.default_model]
            models_caps[provider.default_model] = model_config
            if model_config.get("custom_name"):
                models_names[provider.default_model] = model_config["custom_name"]
    
    return {
        "default": settings.AI_MODEL,
        "models": sorted(list(all_models)),
        "models_caps": models_caps,  # 模型功能信息
        "models_names": models_names,  # 模型自定义显示名称
        "providers": [
            {
                "id": p.id,
                "name": p.name,
                "default_model": p.default_model,
                "models": [m.strip() for m in p.models.split(",") if m.strip()] if p.models else [p.default_model]
            }
            for p in providers
        ]
    }

# ========== 知识库多库管理 + 向量构建接口(新增) ==========

@app.get("/knowledge/bases")
def list_knowledge_bases(db: Session = Depends(get_db)):
    bases = crud.list_knowledge_bases(db)
    return [base.to_dict() for base in bases]

@app.post("/knowledge/bases")
def create_knowledge_base(
    name: str = Form(...),
    description: Optional[str] = Form(None),
    db: Session = Depends(get_db),
):
    kb = crud.create_knowledge_base(db, name=name, description=description)
    return kb.to_dict()

@app.delete("/knowledge/bases/{kb_id}")
def delete_knowledge_base(kb_id: int, db: Session = Depends(get_db)):
    crud.delete_knowledge_base(db, kb_id)
    return {"success": True}

@app.get("/knowledge/documents")
def list_knowledge_documents(
    kb_id: Optional[int] = None,
    db: Session = Depends(get_db),
):
    documents = crud.list_knowledge_documents(db, kb_id=kb_id)
    return [doc.to_dict() for doc in documents]

@app.delete("/knowledge/documents/{doc_id}")
def delete_knowledge_document(doc_id: int, db: Session = Depends(get_db)):
    """删除知识库中的单个文档"""
    crud.delete_knowledge_document(db, doc_id)
    return {"success": True}

@app.post("/knowledge/upload")
def upload_knowledge_file(
    kb_id: Optional[int] = Form(None),
    embedding_model: Optional[str] = Form(None),
    extract_images: Optional[bool] = Form(False),  # 是否提取文档内图片
    vision_model: Optional[str] = Form(None),  # 图片识别用的视觉模型
    file: UploadFile = File(...),
    db: Session = Depends(get_db),
):
    """
    上传一个文件到指定知识库:
    1. 保存原始文件；
    2. 抽取文本(支持 PDF/DOCX/PPTX/XLSX/TXT/MD/CSV/图片 等)；
    3. 如果启用图片提取，识别文档内嵌图片；
    4. 切分为若干段落；
    5. 调用 embedding 接口生成向量；
    6. 存入 KnowledgeDocument + KnowledgeChunk；
    """
    from app.utils.document_parser import extract_text_from_file, get_supported_extensions, set_image_recognition_callback
    
    # 检查文件格式
    ext = os.path.splitext(file.filename)[1].lower()
    supported = get_supported_extensions()
    if ext not in supported:
        raise HTTPException(
            status_code=400, 
            detail=f"不支持的文件格式: {ext}。支持的格式: {', '.join(supported)}"
        )
    
    # 如果需要提取图片或上传的是图片文件，设置图片识别回调
    image_extensions = ['.png', '.jpg', '.jpeg', '.gif', '.bmp', '.webp']
    need_vision = extract_images or ext in image_extensions
    
    if need_vision:
        # 获取视觉模型配置
        selected_vision_model = vision_model
        
        if not selected_vision_model:
            if ext in image_extensions:
                raise HTTPException(
                    status_code=400,
                    detail="上传图片需要配置图片识别方案。请在设置中选择视觉模型。"
                )
            # 如果只是提取文档图片但没配置，则跳过图片提取
            extract_images = False
        else:
            # 设置图片识别回调
            def create_vision_callback(model_name: str):
                # 查找包含该模型的 Provider
                api_base = None
                api_key = None
                
                if model_name.startswith("vision:"):
                    model_name = model_name[7:]
                
                providers = crud.list_providers(db)
                for provider in providers:
                    # 检查 models_config 中是否有该模型
                    if provider.models_config:
                        try:
                            import json as _json
                            config = _json.loads(provider.models_config)
                            if model_name in config:
                                api_base = provider.api_base
                                api_key = provider.api_key
                                break
                        except:
                            pass
                    # 兼容旧的 models 字段
                    if not api_base and provider.models and model_name in provider.models:
                        api_base = provider.api_base
                        api_key = provider.api_key
                        break
                
                if not api_base:
                    # 使用第一个可用的 Provider
                    if providers:
                        api_base = providers[0].api_base
                        api_key = providers[0].api_key
                    else:
                        api_base = settings.AI_API_BASE
                        api_key = settings.AI_API_KEY
                
                def vision_callback(image_bytes: bytes, mime_type: str) -> str:
                    import httpx
                    import base64
                    
                    image_base64 = base64.b64encode(image_bytes).decode('utf-8')
                    image_url = f"data:{mime_type};base64,{image_base64}"
                    
                    headers = {"Content-Type": "application/json"}
                    if api_key:
                        headers["Authorization"] = f"Bearer {api_key}"
                    
                    payload = {
                        "model": model_name,
                        "messages": [
                            {
                                "role": "user",
                                "content": [
                                    {
                                        "type": "text",
                                        "text": "请详细描述这张图片的内容，包括图片中的所有文字、图表、数据、图形等信息。如果有文字，请完整提取出来。"
                                    },
                                    {
                                        "type": "image_url",
                                        "image_url": {"url": image_url}
                                    }
                                ]
                            }
                        ],
                        "max_tokens": 2048
                    }
                    
                    with httpx.Client(timeout=60.0) as client:
                        response = client.post(
                            f"{api_base.rstrip('/')}/chat/completions",
                            headers=headers,
                            json=payload
                        )
                        response.raise_for_status()
                        result = response.json()
                    
                    if "choices" in result and len(result["choices"]) > 0:
                        return result["choices"][0].get("message", {}).get("content", "")
                    return ""
                
                return vision_callback
            
            set_image_recognition_callback(create_vision_callback(selected_vision_model))
    
    # 验证向量模型 - 从 Provider 配置中获取可用的 embedding 模型
    selected_embedding_model = embedding_model or settings.EMBEDDING_MODEL
    
    # 获取所有可用的 embedding 模型(包括 Provider 中配置的)
    available_embedding_models = set(settings.embedding_models)
    all_providers = crud.list_providers(db)
    for provider in all_providers:
        if provider.models_config:
            try:
                import json as _json
                config = _json.loads(provider.models_config)
                for model_name in config.keys():
                    if "embedding" in model_name.lower() or "embed" in model_name.lower():
                        available_embedding_models.add(model_name)
            except:
                pass
    
    if selected_embedding_model and selected_embedding_model not in available_embedding_models:
        selected_embedding_model = None
    
    # 1. 保存文件
    kb_dir = os.path.join(UPLOAD_DIR, "knowledge")
    os.makedirs(kb_dir, exist_ok=True)
    save_path = os.path.join(kb_dir, file.filename)
    with open(save_path, "wb") as f:
        shutil.copyfileobj(file.file, f)

    # 2. 提取文本(支持多种格式，可选提取图片)
    try:
        content = extract_text_from_file(save_path, extract_images=extract_images)
    except ImportError as e:
        raise HTTPException(status_code=500, detail=f"缺少依赖库: {e}")
    except ValueError as e:
        raise HTTPException(status_code=400, detail=str(e))
    except Exception as e:
        raise HTTPException(status_code=400, detail=f"文件解析失败: {e}")
    finally:
        # 清理回调
        set_image_recognition_callback(None)

    # 3. 智能切分段落(chunk_size = 512)
    CHUNK_SIZE = 512
    CHUNK_OVERLAP = 50  # 重叠部分，保持上下文连贯
    
    paragraphs: List[str] = []
    current_chunk = ""
    
    for line in content.splitlines():
        line = line.strip()
        if not line:
            continue
        
        # 如果当前块加上新行不超过 chunk_size，则合并
        if len(current_chunk) + len(line) + 1 <= CHUNK_SIZE:
            if current_chunk:
                current_chunk += "\n" + line
            else:
                current_chunk = line
        else:
            # 保存当前块
            if current_chunk:
                paragraphs.append(current_chunk)
            
            # 如果单行超过 chunk_size，按句子切分
            if len(line) > CHUNK_SIZE:
                import re
                sentences = re.split(r'(?<=[。！？!?.;])\s*', line)
                temp_chunk = ""
                for sent in sentences:
                    sent = sent.strip()
                    if not sent:
                        continue
                    if len(temp_chunk) + len(sent) + 1 <= CHUNK_SIZE:
                        if temp_chunk:
                            temp_chunk += sent
                        else:
                            temp_chunk = sent
                    else:
                        if temp_chunk:
                            paragraphs.append(temp_chunk)
                        temp_chunk = sent
                current_chunk = temp_chunk
            else:
                current_chunk = line
    
    # 保存最后一个块
    if current_chunk:
        paragraphs.append(current_chunk)

    if not paragraphs:
        raise HTTPException(status_code=400, detail="文件中未检测到有效文本内容。")

    # 4. 生成向量(如果有向量模型)
    embeddings = None
    if selected_embedding_model:
        try:
            # 配置 ai_manager 使用第一个可用的 Provider
            all_providers = crud.list_providers(db)
            if all_providers:
                provider = all_providers[0]
                ai_manager.set_provider(
                    api_base=provider.api_base,
                    api_key=provider.api_key,
                    default_model=provider.default_model,
                )
            
            embeddings = ai_manager.create_embedding(paragraphs, model=selected_embedding_model)
        except Exception as e:
            chat_logger.error(f"向量生成失败: {e}")
            # 删除已保存的文件
            try:
                os.remove(save_path)
            except:
                pass
            raise HTTPException(status_code=500, detail=f"向量生成失败: {str(e)}，文件未加入知识库。")
    
    # 如果选择了向量模型但没有生成向量，不写入数据库
    if selected_embedding_model and (not embeddings or len(embeddings) != len(paragraphs)):
        # 删除已保存的文件
        try:
            os.remove(save_path)
        except:
            pass
        raise HTTPException(status_code=500, detail="向量生成失败或数量不匹配，文件未加入知识库。")
    
    # 如果没有选择向量模型，也不写入数据库
    if not selected_embedding_model:
        # 删除已保存的文件
        try:
            os.remove(save_path)
        except:
            pass
        raise HTTPException(status_code=400, detail="请选择向量模型，否则文件无法用于知识库搜索。")

    # 5. 写入 DB(只有成功生成向量才写入)
    doc = crud.create_knowledge_document(
        db,
        kb_id=kb_id,
        file_name=file.filename,
        file_path=save_path,
        content=content[:2000],
        embedding_model=selected_embedding_model,
    )

    # 存储向量块
    chunks_data = []
    for idx, (para, emb) in enumerate(zip(paragraphs, embeddings)):
        chunks_data.append((idx, para, emb))
    crud.create_knowledge_chunks(db, document_id=doc.id, chunks=chunks_data)

    return {
        "success": True, 
        "document": doc.to_dict(),
        "chunks_count": len(paragraphs),
    }

# ========== MCP 服务器管理接口 ==========

@app.get("/mcp/servers")
async def get_mcp_servers(db: Session = Depends(get_db)):
    """获取 MCP 服务器列表"""
    # 从数据库获取配置
    saved_config = crud.get_setting(db, "mcp_servers")
    servers_config = json.loads(saved_config.value) if saved_config else []
    
    # 获取运行状态
    result = []
    for config in servers_config:
        server_name = config.get("name", "")
        is_running = False
        tools = []
        
        if server_name in mcp_client.servers:
            server = mcp_client.servers[server_name]
            is_running = server.process is not None and server.process.poll() is None
            tools = [{"name": t.name, "description": t.description} for t in server.tools]
        
        result.append({
            "name": server_name,
            "type": config.get("type", "stdio"),
            "command": config.get("command", ""),
            "args": config.get("args", []),
            "url": config.get("url", ""),
            "env": config.get("env", {}),
            "enabled": config.get("enabled", True),
            "running": is_running,
            "tools": tools
        })
    
    return {"servers": result}

@app.post("/mcp/servers/test")
async def test_mcp_server(
    name: str = Form(...),
    type: str = Form("stdio"),
    command: str = Form(""),
    args: str = Form(""),
    url: str = Form(""),
    env: str = Form(""),
):
    """测试 MCP 服务器连接(不保存到数据库)"""
    if type != "stdio":
        return {"success": False, "error": "目前只支持 stdio 类型"}
    
    if not command:
        return {"success": False, "error": "请填写命令"}
    
    # 解析参数
    args_list = [a.strip() for a in args.split() if a.strip()] if args else []
    
    # 解析环境变量
    env_dict = {}
    if env:
        for line in env.strip().split("\n"):
            if "=" in line:
                k, v = line.split("=", 1)
                env_dict[k.strip()] = v.strip()
    
    # 创建临时客户端测试
    from app.ai.mcp_client import MCPClient
    test_client = MCPClient()
    test_client.add_server("_test_", command, args_list, env_dict)
    
    try:
        success = await test_client.start_server("_test_")
        if success:
            server = test_client.servers["_test_"]
            tools = [{"name": t.name, "description": t.description} for t in server.tools]
            await test_client.stop_server("_test_")
            return {"success": True, "tools": tools}
        else:
            return {"success": False, "error": "无法启动服务器"}
    except Exception as e:
        await test_client.stop_all()
        return {"success": False, "error": str(e)}

@app.post("/mcp/servers")
async def add_mcp_server(
    name: str = Form(...),
    type: str = Form("stdio"),
    command: str = Form(""),
    args: str = Form(""),
    url: str = Form(""),
    env: str = Form(""),
    enabled: bool = Form(True),
    db: Session = Depends(get_db)
):
    """添加或更新 MCP 服务器"""
    # 解析参数
    args_list = [a.strip() for a in args.split() if a.strip()] if args else []
    
    # 解析环境变量
    env_dict = {}
    if env:
        for line in env.strip().split("\n"):
            if "=" in line:
                k, v = line.split("=", 1)
                env_dict[k.strip()] = v.strip()
    
    # 获取现有配置
    saved_config = crud.get_setting(db, "mcp_servers")
    servers_config = json.loads(saved_config.value) if saved_config else []
    
    # 检查是否已存在(更新)
    existing_idx = None
    for i, config in enumerate(servers_config):
        if config.get("name") == name:
            existing_idx = i
            break
    
    # 新配置
    new_config = {
        "name": name,
        "type": type,
        "command": command,
        "args": args_list,
        "url": url,
        "env": env_dict,
        "enabled": enabled
    }
    
    if existing_idx is not None:
        servers_config[existing_idx] = new_config
    else:
        servers_config.append(new_config)
    
    # 保存到数据库
    try:
        crud.set_setting(db, "mcp_servers", json.dumps(servers_config, ensure_ascii=False))
    except Exception as e:
        return {"success": False, "error": f"保存失败: {e}"}
    
    # 更新客户端配置
    if name in mcp_client.servers:
        await mcp_client.stop_server(name)
        del mcp_client.servers[name]
    
    if enabled and type == "stdio" and command:
        mcp_client.add_server(name, command, args_list, env_dict)
        try:
            await mcp_client.start_server(name)
            return {"success": True, "message": f"MCP Server {name} 已保存并启动"}
        except Exception as e:
            return {"success": True, "message": f"MCP Server {name} 已保存，但启动失败: {e}"}
    
    return {"success": True, "message": f"MCP Server {name} 已保存"}

@app.delete("/mcp/servers/{name}")
async def delete_mcp_server(name: str, db: Session = Depends(get_db)):
    """删除 MCP 服务器"""
    # 停止服务器
    await mcp_client.stop_server(name)
    if name in mcp_client.servers:
        del mcp_client.servers[name]
    
    # 从数据库删除
    saved_config = crud.get_setting(db, "mcp_servers")
    servers_config = json.loads(saved_config.value) if saved_config else []
    servers_config = [c for c in servers_config if c.get("name") != name]
    crud.set_setting(db, "mcp_servers", json.dumps(servers_config, ensure_ascii=False))
    
    return {"success": True, "message": f"服务器 {name} 已删除"}

@app.post("/mcp/servers/{name}/start")
async def start_mcp_server(name: str, db: Session = Depends(get_db)):
    """启动 MCP 服务器"""
    # 获取配置
    saved_config = crud.get_setting(db, "mcp_servers")
    servers_config = json.loads(saved_config.value) if saved_config else []
    
    config = None
    for c in servers_config:
        if c.get("name") == name:
            config = c
            break
    
    if not config:
        return {"success": False, "error": f"服务器 {name} 不存在"}
    
    # 添加到客户端(如果还没有)
    if name not in mcp_client.servers:
        mcp_client.add_server(name, config["command"], config.get("args", []))
    
    # 启动
    try:
        success = await mcp_client.start_server(name)
        if success:
            server = mcp_client.servers[name]
            tools = [{"name": t.name, "description": t.description} for t in server.tools]
            return {"success": True, "message": f"服务器 {name} 已启动", "tools": tools}
        else:
            return {"success": False, "error": f"服务器 {name} 启动失败"}

    except Exception as e:
        return {"success": False, "error": str(e)}

@app.post("/mcp/servers/{name}/stop")
async def stop_mcp_server(name: str):
    """停止 MCP 服务器"""
    await mcp_client.stop_server(name)
    return {"success": True, "message": f"服务器 {name} 已停止"}

@app.get("/mcp/tools")
async def get_mcp_tools():
    """获取所有可用的 MCP 工具"""
    tools = mcp_client.get_tools_for_display()
    return {"tools": tools}

@app.post("/mcp/call")
async def call_mcp_tool(
    server: str = Form(...),
    tool: str = Form(...),
    arguments: str = Form("{}")
):
    """直接调用 MCP 工具(用于测试)"""
    try:
        args = json.loads(arguments)
    except:
        args = {}
    
    result = await mcp_client.call_tool(server, tool, args)
    return result

# ========== 系统设置接口(新增) ==========

@app.get("/settings")
def get_settings(db: Session = Depends(get_db)):
    """获取系统设置"""
    # 从数据库获取保存的设置
    saved_settings = crud.get_all_settings(db)
    settings_dict = {s.key: s.value for s in saved_settings}
    
    # 对于敏感信息，只返回掩码
    tavily_key = settings_dict.get("tavily_api_key", "")
    tavily_key_masked = ""
    if tavily_key:
        if len(tavily_key) > 4:
            tavily_key_masked = "***" + tavily_key[-4:]
        else:
            tavily_key_masked = "****"
    
    return {
        "layout_scale": settings_dict.get("layout_scale", "normal"),
        "auto_title_model": settings_dict.get("auto_title_model", "current"),
        "default_vision_model": settings_dict.get("default_vision_model", ""),
        "default_search_source": settings_dict.get("default_search_source", "duckduckgo"),
        "tavily_api_key": tavily_key_masked,  # 返回掩码而非完整key
        "global_api_key": settings_dict.get("global_api_key", getattr(settings, 'AI_API_KEY', '')),
        "global_api_base": settings_dict.get("global_api_base", getattr(settings, 'AI_API_BASE', 'https://api.openai.com/v1')),
        "global_default_model": settings_dict.get("global_default_model", getattr(settings, 'AI_MODEL', 'gpt-4o-mini')),
    }

@app.post("/settings")
def update_settings(
    layout_scale: Optional[str] = Form(None),
    auto_title_model: Optional[str] = Form(None),
    default_vision_model: Optional[str] = Form(None),
    theme: Optional[str] = Form(None),
    language: Optional[str] = Form(None),
    default_search_source: Optional[str] = Form(None),
    tavily_api_key: Optional[str] = Form(None),
    global_api_key: Optional[str] = Form(None),
    global_api_base: Optional[str] = Form(None),
    global_default_model: Optional[str] = Form(None),
    db: Session = Depends(get_db),
):
    """更新系统设置"""
    settings_data = {}
    
    # 保存设置到数据库
    if layout_scale:
        crud.set_setting(db, "layout_scale", layout_scale)
        settings_data["layout_scale"] = layout_scale
    if auto_title_model:
        crud.set_setting(db, "auto_title_model", auto_title_model)
        settings_data["auto_title_model"] = auto_title_model
    if default_vision_model is not None:  # 允许空字符串(表示不启用)
        crud.set_setting(db, "default_vision_model", default_vision_model)
        settings_data["default_vision_model"] = default_vision_model
    if theme:
        crud.set_setting(db, "theme", theme)
        settings_data["theme"] = theme
    if language:
        crud.set_setting(db, "language", language)
        settings_data["language"] = language
    if default_search_source:
        crud.set_setting(db, "default_search_source", default_search_source)
        settings_data["default_search_source"] = default_search_source
    if tavily_api_key is not None:  # 允许空字符串
        crud.set_setting(db, "tavily_api_key", tavily_api_key)
        settings_data["tavily_api_key"] = tavily_api_key
    
    # 新增:全局API配置
    if global_api_key is not None:
        crud.set_setting(db, "global_api_key", global_api_key)
        settings_data["global_api_key"] = global_api_key
        # 同时更新AI管理器的配置
        ai_manager._provider.api_key = global_api_key
        # 更新环境变量(如果需要持久化)
        import os
        os.environ["AI_API_KEY"] = global_api_key
        
    if global_api_base is not None:
        crud.set_setting(db, "global_api_base", global_api_base)
        settings_data["global_api_base"] = global_api_base
        ai_manager._provider.api_base = global_api_base.rstrip("/")
        os.environ["AI_API_BASE"] = global_api_base
        
    if global_default_model is not None:
        crud.set_setting(db, "global_default_model", global_default_model)
        settings_data["global_default_model"] = global_default_model
        ai_manager._provider.default_model = global_default_model
        os.environ["AI_MODEL"] = global_default_model
    
    return {"success": True, "settings": settings_data}

@app.get("/logs/export")
def export_logs(hours: int = 24):
    """导出日志文件"""
    import zipfile
    import io
    from datetime import datetime, timedelta
    from pathlib import Path
    
    logs_dir = Path("logs")
    if not logs_dir.exists():
        raise HTTPException(status_code=404, detail="日志目录不存在")
    
    cutoff_time = datetime.now() - timedelta(hours=hours)
    
    # 创建内存中的 ZIP 文件
    zip_buffer = io.BytesIO()
    
    with zipfile.ZipFile(zip_buffer, 'w', zipfile.ZIP_DEFLATED) as zipf:
        # 添加系统信息
        import platform
        system_info = {
            "timestamp": datetime.now().isoformat(),
            "platform": platform.platform(),
            "python_version": platform.python_version(),
            "hours_collected": hours
        }
        zipf.writestr("system_info.json", json.dumps(system_info, indent=2, ensure_ascii=False))
        
        # 收集日志文件
        log_files = ["main.log", "api.log", "chat.log", "token.log", "database.log", "error.log"]
        collected_count = 0
        
        for log_file in log_files:
            log_path = logs_dir / log_file
            if not log_path.exists():
                continue
            
            try:
                with open(log_path, 'r', encoding='utf-8') as f:
                    lines = f.readlines()
                
                # 过滤最近的日志
                recent_lines = []
                for line in lines:
                    try:
                        if len(line) > 19:
                            timestamp_str = line[:19]
                            log_time = datetime.strptime(timestamp_str, "%Y-%m-%d %H:%M:%S")
                            if log_time >= cutoff_time:
                                recent_lines.append(line)
                    except:
                        recent_lines.append(line)
                
                if recent_lines:
                    zipf.writestr(f"logs/{log_file}", ''.join(recent_lines))
                    collected_count += 1
            except Exception:
                pass
        
        # 添加说明文件
        readme = f"""日志导出
生成时间: {datetime.now().isoformat()}
收集范围: 最近 {hours} 小时
文件数量: {collected_count}
"""
        zipf.writestr("README.txt", readme)
    
    zip_buffer.seek(0)
    
    # 生成文件名
    timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
    filename = f"debug_logs_{timestamp}.zip"
    
    from fastapi.responses import StreamingResponse
    return StreamingResponse(
        zip_buffer,
        media_type="application/zip",
        headers={"Content-Disposition": f"attachment; filename={filename}"}
    )

@app.post("/search/test")
def test_search_connection(
    source: str = Form(...),
    query: str = Form("test search"),
    tavily_api_key: Optional[str] = Form(None),
    db: Session = Depends(get_db),
):
    """测试搜索API连接"""
    try:
        if source == "duckduckgo":
            # DuckDuckGo 不需要 API Key，直接测试
            import requests
            params = {"q": query, "format": "json"}
            response = requests.get(
                "https://api.duckduckgo.com/",
                params=params,
                timeout=10
            )
            response.raise_for_status()
            
        elif source == "tavily":
            if not tavily_api_key:
                setting = crud.get_setting(db, "tavily_api_key")
                tavily_api_key = setting.value if setting else None
                if not tavily_api_key:
                    raise HTTPException(status_code=400, detail="请提供Tavily API Key")
            
            import requests
            payload = {
                "api_key": tavily_api_key,
                "query": query,
                "max_results": 1
            }

            response = requests.post(
                "https://api.tavily.com/search",
                json=payload,
                timeout=10
            )
            response.raise_for_status()
            
        else:
            raise HTTPException(status_code=400, detail="不支持的搜索源")
        
        return {"success": True, "message": f"{source.title()} 搜索连接测试成功"}
        
    except requests.exceptions.RequestException as e:
        raise HTTPException(status_code=500, detail=f"搜索API连接失败: {str(e)}")
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"测试失败: {str(e)}")

@app.post("/test-api-connection")
def test_api_connection():
    """测试全局API连接"""
    try:
        if not ai_manager.is_configured():
            return {"success": False, "error": "请先配置API Key和API Base"}
        
        # 发送测试请求
        test_messages = [{"role": "user", "content": "Hello"}]
        result = ai_manager.chat(test_messages, stream=False)
        
        return {"success": True, "message": "API连接测试成功"}
    except Exception as e:
        return {"success": False, "error": f"API连接测试失败: {str(e)}"}

@app.post("/test-provider-connection")
def test_provider_connection(
    api_base: str = Form(...),
    api_key: str = Form(...),
    model: str = Form(...),
):
    """测试指定Provider的连接"""
    try:
        # 创建临时的AI管理器实例进行测试
        temp_manager = AIManager()
        temp_manager.set_provider(
            api_base=api_base,
            api_key=api_key,
            default_model=model
        )
        
        # 发送测试请求
        test_messages = [{"role": "user", "content": "Hello"}]
        result = temp_manager.chat(test_messages, stream=False)
        
        return {"success": True, "message": "Provider连接测试成功"}
    except Exception as e:
        return {"success": False, "error": f"Provider连接测试失败: {str(e)}"}

@app.get("/api-status")
def get_api_status():
    """获取API配置状态"""
    return {
        "configured": ai_manager.is_configured(),
        "api_base": ai_manager._provider.api_base,
        "has_api_key": bool(ai_manager._provider.api_key),
        "default_model": ai_manager._provider.default_model
    }

# ========== 知识图谱接口 ==========

@app.get("/knowledge/graph/stats")
def get_knowledge_graph_stats(
    kb_id: Optional[int] = None,
    db: Session = Depends(get_db),
):
    """获取知识图谱统计信息"""
    stats = crud.get_knowledge_graph_stats(db, kb_id)
    return stats

@app.get("/knowledge/graph/entities")
def list_knowledge_entities(
    kb_id: Optional[int] = None,
    entity_type: Optional[str] = None,
    limit: int = 100,
    db: Session = Depends(get_db),
):
    """列出知识图谱实体"""
    entities = crud.list_entities(db, kb_id=kb_id, entity_type=entity_type, limit=limit)
    return [e.to_dict() for e in entities]

@app.get("/knowledge/graph/entities/search")
def search_knowledge_entities(
    query: str,
    kb_id: Optional[int] = None,
    limit: int = 10,
    db: Session = Depends(get_db),
):
    """搜索知识图谱实体"""
    entities = crud.search_entities(db, query, kb_id=kb_id, limit=limit)
    return [e.to_dict() for e in entities]

@app.get("/knowledge/graph/entities/{entity_id}")
def get_knowledge_entity(
    entity_id: int,
    db: Session = Depends(get_db),
):
    """获取单个实体详情"""
    entity = crud.get_entity(db, entity_id)
    if not entity:
        raise HTTPException(status_code=404, detail="实体不存在")
    return entity.to_dict()

@app.get("/knowledge/graph/entities/{entity_id}/relations")
def get_entity_relations(
    entity_id: int,
    max_depth: int = 2,
    db: Session = Depends(get_db),
):
    """获取实体的关系网络"""
    entity = crud.get_entity(db, entity_id)
    if not entity:
        raise HTTPException(status_code=404, detail="实体不存在")
    
    related = crud.get_related_entities(db, entity_id, max_depth=max_depth)
    
    return {
        "entity": entity.to_dict(),
        "related": [
            {
                "entity": item["entity"].to_dict(),
                "relation": item["relation"].to_dict(),
                "depth": item["depth"],
            }
            for item in related
        ]
    }

@app.get("/knowledge/graph/relations")
def list_knowledge_relations(
    kb_id: Optional[int] = None,
    entity_id: Optional[int] = None,
    relation_type: Optional[str] = None,
    limit: int = 100,
    db: Session = Depends(get_db),
):
    """列出知识图谱关系"""
    relations = crud.list_relations(
        db, 
        kb_id=kb_id, 
        entity_id=entity_id, 
        relation_type=relation_type, 
        limit=limit
    )
    return [r.to_dict() for r in relations]

@app.post("/knowledge/graph/entities")
def create_knowledge_entity(
    kb_id: Optional[int] = Form(None),
    name: str = Form(...),
    entity_type: str = Form(...),
    description: Optional[str] = Form(None),
    db: Session = Depends(get_db),
):
    """手动创建实体"""
    entity = crud.create_entity(
        db,
        kb_id=kb_id,
        name=name,
        entity_type=entity_type,
        description=description,
    )
    return entity.to_dict()

@app.post("/knowledge/graph/relations")
def create_knowledge_relation(
    kb_id: Optional[int] = Form(None),
    source_id: int = Form(...),
    target_id: int = Form(...),
    relation_type: str = Form(...),
    description: Optional[str] = Form(None),
    db: Session = Depends(get_db),
):
    """手动创建关系"""
    # 验证实体存在
    source = crud.get_entity(db, source_id)
    target = crud.get_entity(db, target_id)
    if not source or not target:
        raise HTTPException(status_code=404, detail="源实体或目标实体不存在")
    
    relation = crud.create_relation(
        db,
        kb_id=kb_id,
        source_id=source_id,
        target_id=target_id,
        relation_type=relation_type,
        description=description,
    )
    return relation.to_dict()

@app.delete("/knowledge/graph/entities/{entity_id}")
def delete_knowledge_entity(entity_id: int, db: Session = Depends(get_db)):
    """删除实体(会级联删除相关关系)"""
    crud.delete_entity(db, entity_id)
    return {"success": True}

@app.delete("/knowledge/graph/relations/{relation_id}")
def delete_knowledge_relation(relation_id: int, db: Session = Depends(get_db)):
    """删除关系"""
    crud.delete_relation(db, relation_id)
    return {"success": True}

@app.get("/knowledge/graph/context")
def get_graph_context(
    query: str,
    kb_id: Optional[int] = None,
    max_entities: int = 5,
    db: Session = Depends(get_db),
):
    """基于查询获取知识图谱上下文(用于增强 RAG)"""
    from app.ai.knowledge_graph import search_graph_context
    context = search_graph_context(db, query, kb_id=kb_id, max_entities=max_entities)
    return {"context": context}

@app.get("/knowledge/embedding-models")
def get_embedding_models(db: Session = Depends(get_db)):
    """获取可用的向量模型列表 - 基于用户配置的Provider，按Provider分组"""
    import json
    
    # 获取所有Provider
    providers = crud.list_providers(db)
    
    # 收集所有Provider中的向量模型，按Provider分组
    embedding_models = []
    all_models = set()
    models_names = {}  # 自定义名称
    
    # 添加全局默认向量模型
    for model in settings.embedding_models:
        if model not in all_models:
            embedding_models.append({
                "model": model,
                "provider_id": None,
                "provider_name": "默认",
                "custom_name": None
            })
            all_models.add(model)
    
    # 从Provider中提取向量模型
    for provider in providers:
        if provider.models:
            # 解析模型配置获取自定义名称
            config = {}
            if provider.models_config:
                try:
                    config = json.loads(provider.models_config)
                except:
                    pass
            
            provider_models = [m.strip() for m in provider.models.split(",") if m.strip()]
            # 过滤出向量模型(通常包含embedding关键字)
            for model in provider_models:
                if "embedding" in model.lower() or "embed" in model.lower():
                    custom_name = config.get(model, {}).get("custom_name") if config.get(model) else None
                    embedding_models.append({
                        "model": model,
                        "provider_id": provider.id,
                        "provider_name": provider.name,
                        "custom_name": custom_name
                    })
                    if custom_name:
                        models_names[model] = custom_name
                    all_models.add(model)
    
    # 如果没有找到任何向量模型，返回空列表
    if not embedding_models:
        return {
            "default": None,
            "models": [],
            "models_by_provider": [],
            "models_names": {},
            "message": "当前Provider中未配置向量模型，请在Provider设置中添加向量模型"
        }
    
    return {
        "default": settings.EMBEDDING_MODEL if settings.EMBEDDING_MODEL in all_models else embedding_models[0]["model"],
        "models": sorted(list(all_models)),
        "models_by_provider": embedding_models,
        "models_names": models_names,
    }

@app.get("/models/vision")
def get_vision_models(db: Session = Depends(get_db)):
    """获取可用的视觉模型列表 - 基于 models_config 中的 vision 标记"""
    import json
    
    # 获取所有Provider
    providers = crud.list_providers(db)
    
    # 收集所有Provider中的视觉模型
    vision_models = []
    all_models = set()
    models_names = {}
    
    # 从Provider 的 models_config 中提取标记了 vision 能力的模型
    for provider in providers:
        if provider.models_config:
            try:
                config = json.loads(provider.models_config)
                for model_name, caps in config.items():
                    if caps.get("vision"):
                        if model_name not in all_models:
                            custom_name = caps.get("custom_name")
                            vision_models.append({
                                "model": model_name,
                                "provider_id": provider.id,
                                "provider_name": provider.name,
                                "custom_name": custom_name
                            })
                            all_models.add(model_name)
                            if custom_name:
                                models_names[model_name] = custom_name
            except:
                pass
    
    return {
        "default": vision_models[0]["model"] if vision_models else None,
        "models": sorted(list(all_models)),
        "models_by_provider": vision_models,
        "models_names": models_names,
    }

@app.get("/models/rerank")
def get_rerank_models(db: Session = Depends(get_db)):
    """获取可用的重排模型列表 - 按Provider分组"""
    import json
    
    # 获取所有Provider
    providers = crud.list_providers(db)
    
    # 收集所有Provider中的重排模型，按Provider分组
    rerank_models = []
    all_models = set()
    models_names = {}  # 自定义名称
    
    # 从Provider 的 models_config 中提取重排模型(模型名包含 rerank)
    for provider in providers:
        if provider.models_config:
            try:
                config = json.loads(provider.models_config)
                for model_name, caps in config.items():
                    if "rerank" in model_name.lower():
                        if model_name not in all_models:
                            custom_name = caps.get("custom_name") if isinstance(caps, dict) else None
                            rerank_models.append({
                                "model": model_name,
                                "provider_id": provider.id,
                                "provider_name": provider.name,
                                "custom_name": custom_name
                            })
                            all_models.add(model_name)
                            if custom_name:
                                models_names[model_name] = custom_name
            except:
                pass
    
    return {
        "default": rerank_models[0]["model"] if rerank_models else None,
        "models": sorted(list(all_models)),
        "models_by_provider": rerank_models,
        "models_names": models_names,
    }

@app.post("/images/generate")
@log_api_call
def generate_image(
    prompt: str = Form(...),
    model: str = Form(...),
    size: str = Form("1024x1024"),
    n: int = Form(1),
    provider_id: Optional[int] = Form(None),
    conversation_id: Optional[int] = Form(None),
    db: Session = Depends(get_db),
):
    """
    生成图像 API
    
    参数：
    - prompt: 图像描述
    - model: 生图模型名称
    - size: 图像尺寸
    - n: 生成数量
    - provider_id: 使用的 Provider ID
    - conversation_id: 关联的对话 ID(可选，用于保存到对话历史)
    """
    from app.ai.ai_manager import AIManager
    
    # 配置 Provider
    ai = AIManager()
    
    if provider_id:
        provider = crud.get_provider(db, provider_id)
        if provider:
            ai.set_provider(
                api_base=provider.api_base,
                api_key=provider.api_key,
                default_model=provider.default_model
            )
    
    # 调用生图 API
    result = ai.generate_image(
        prompt=prompt,
        model=model,
        size=size,
        n=n,
        response_format="url"  # 优先返回 URL
    )
    
    if not result["success"]:
        # 如果 URL 格式失败，尝试 base64 格式
        result = ai.generate_image(
            prompt=prompt,
            model=model,
            size=size,
            n=n,
            response_format="b64_json"
        )
    
    # 如果指定了对话 ID，将生成的图片保存到对话历史
    if conversation_id and result["success"] and result["images"]:
        conversation = crud.get_conversation(db, conversation_id)
        if conversation:
            # 构建包含图片的消息内容
            image_content = f"**生成图片** (模型: {model}, 尺寸: {size})\n\n"
            for i, img in enumerate(result["images"]):
                if "url" in img:
                    image_content += f"![生成的图片 {i+1}]({img['url']})\n\n"
                elif "b64_json" in img:
                    image_content += f"![生成的图片 {i+1}](data:image/png;base64,{img['b64_json']})\n\n"
            
            # 保存用户的生图请求
            crud.create_message(db, conversation_id, "user", f"[生图请求] {prompt}")
            # 保存生成结果
            crud.create_message(db, conversation_id, "assistant", image_content)
    
    return result

@app.get("/models/image-gen")
def get_image_gen_models(db: Session = Depends(get_db)):
    """获取可用的生图模型列表"""
    providers = crud.list_providers(db)
    
    image_gen_models = []
    
    for provider in providers:
        if provider.models_config:
            try:
                import json
                models_config = json.loads(provider.models_config)
                for model_name, caps in models_config.items():
                    if caps.get("image_gen"):
                        image_gen_models.append({
                            "model": model_name,
                            "provider_id": provider.id,
                            "provider_name": provider.name,
                            "custom_name": caps.get("custom_name", "")
                        })
            except:
                pass
    
    return {
        "models": image_gen_models
    }

# ========== 静态页面 ==========

@app.get("/")
def index():
    # 返回新的分离后的前端页面
    from fastapi.responses import HTMLResponse

    frontend_path = os.path.join(os.path.dirname(__file__), "..", "frontend", "index_new.html")
    frontend_path = os.path.abspath(frontend_path)
    if not os.path.exists(frontend_path):
        # 如果新文件不存在，回退到原文件
        frontend_path = os.path.join(os.path.dirname(__file__), "..", "frontend", "index.html")
        frontend_path = os.path.abspath(frontend_path)
        if not os.path.exists(frontend_path):
            return HTMLResponse("<h1>Frontend not found</h1>", status_code=404)
    
    with open(frontend_path, "r", encoding="utf-8") as f:
        html = f.read()
    return HTMLResponse(html)

@app.get("/style.css")
def get_css():
    # 返回CSS文件
    frontend_path = os.path.join(os.path.dirname(__file__), "..", "frontend", "style.css")
    frontend_path = os.path.abspath(frontend_path)
    if not os.path.exists(frontend_path):
        raise HTTPException(status_code=404, detail="CSS file not found")
    return FileResponse(frontend_path, media_type="text/css")

@app.get("/script.js")
def get_js():
    # 返回JavaScript文件
    frontend_path = os.path.join(os.path.dirname(__file__), "..", "frontend", "script.js")
    frontend_path = os.path.abspath(frontend_path)
    if not os.path.exists(frontend_path):
        raise HTTPException(status_code=404, detail="JavaScript file not found")
    return FileResponse(frontend_path, media_type="application/javascript")

@app.get("/markdown.js")
def get_markdown_js():
    # 返回Markdown渲染JavaScript文件
    frontend_path = os.path.join(os.path.dirname(__file__), "..", "frontend", "markdown.js")
    frontend_path = os.path.abspath(frontend_path)
    if not os.path.exists(frontend_path):
        raise HTTPException(status_code=404, detail="Markdown JavaScript file not found")
    return FileResponse(frontend_path, media_type="application/javascript")

@app.get("/favicon.ico")
def get_favicon():
    # 返回favicon文件
    frontend_path = os.path.join(os.path.dirname(__file__), "..", "frontend", "favicon.ico")
    frontend_path = os.path.abspath(frontend_path)
    if not os.path.exists(frontend_path):
        # 如果没有favicon文件，返回一个简单的响应
        from fastapi.responses import Response
        return Response(content="", media_type="image/x-icon")
    return FileResponse(frontend_path, media_type="image/x-icon")

@app.get("/lib/{filename:path}")
def get_lib_file(filename: str):
    """返回lib目录下的静态文件(JS库、CSS等)"""
    frontend_path = os.path.join(os.path.dirname(__file__), "..", "frontend", "lib", filename)
    frontend_path = os.path.abspath(frontend_path)
    
    # 安全检查:确保路径在lib目录内
    lib_dir = os.path.abspath(os.path.join(os.path.dirname(__file__), "..", "frontend", "lib"))
    if not frontend_path.startswith(lib_dir):
        raise HTTPException(status_code=403, detail="Access denied")
    
    if not os.path.exists(frontend_path):
        raise HTTPException(status_code=404, detail=f"Library file not found: {filename}")
    
    # 根据文件扩展名设置MIME类型
    ext = os.path.splitext(filename)[1].lower()
    mime_types = {
        ".js": "application/javascript",
        ".css": "text/css",
        ".json": "application/json",
        ".map": "application/json",
        ".woff": "font/woff",
        ".woff2": "font/woff2",
        ".ttf": "font/ttf",
        ".eot": "application/vnd.ms-fontobject",
    }
    media_type = mime_types.get(ext, "application/octet-stream")
    
    return FileResponse(frontend_path, media_type=media_type)

@app.get("/render-logger.js")
def get_render_logger_js():
    # 返回渲染日志JavaScript文件
    frontend_path = os.path.join(os.path.dirname(__file__), "..", "frontend", "render-logger.js")
    frontend_path = os.path.abspath(frontend_path)
    if not os.path.exists(frontend_path):
        raise HTTPException(status_code=404, detail="Render logger JavaScript file not found")
    return FileResponse(frontend_path, media_type="application/javascript")

# 前端日志接收API
@app.post("/api/frontend-log")
async def receive_frontend_log(log_entry: dict):
    """接收前端日志并写入文件"""
    import datetime
    log_dir = os.path.join(os.path.dirname(__file__), "..", "logs")
    log_file = os.path.join(log_dir, "frontend-render.log")
    
    try:
        os.makedirs(log_dir, exist_ok=True)
        with open(log_file, "a", encoding="utf-8") as f:
            timestamp = log_entry.get("timestamp", datetime.datetime.now().isoformat())
            level = log_entry.get("level", "info").upper()
            category = log_entry.get("category", "UNKNOWN")
            message = log_entry.get("message", "")
            data = log_entry.get("data", "")
            session_id = log_entry.get("sessionId", "")
            
            log_line = f"{timestamp} [{level}][{category}][{session_id}] {message}"
            if data:
                log_line += f" | {data}"
            f.write(log_line + "\n")
        return {"status": "ok"}
    except Exception as e:
        return {"status": "error", "message": str(e)}

@app.post("/api/frontend-log/batch")
async def receive_frontend_logs_batch(data: dict):
    """批量接收前端日志"""
    import datetime
    log_dir = os.path.join(os.path.dirname(__file__), "..", "logs")
    log_file = os.path.join(log_dir, "frontend-render.log")
    
    logs = data.get("logs", [])
    try:
        os.makedirs(log_dir, exist_ok=True)
        with open(log_file, "a", encoding="utf-8") as f:
            for log_entry in logs:
                timestamp = log_entry.get("timestamp", datetime.datetime.now().isoformat())
                level = log_entry.get("level", "info").upper()
                category = log_entry.get("category", "UNKNOWN")
                message = log_entry.get("message", "")
                log_data = log_entry.get("data", "")
                session_id = log_entry.get("sessionId", "")
                
                log_line = f"{timestamp} [{level}][{category}][{session_id}] {message}"
                if log_data:
                    log_line += f" | {log_data}"
                f.write(log_line + "\n")
        return {"status": "ok", "count": len(logs)}
    except Exception as e:
        return {"status": "error", "message": str(e)}
