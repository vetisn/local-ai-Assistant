"""
文档解析工具 - 支持多种文件格式
支持: PDF, DOCX, DOC, PPTX, XLSX, TXT, MD, CSV, 图片
支持提取文档内嵌图片并用视觉模型识别
"""
import os
import base64
import io
from typing import Optional, List, Tuple, Callable


# 图片识别回调函数（由外部设置）
_image_recognition_callback: Optional[Callable[[bytes, str], str]] = None


def set_image_recognition_callback(callback: Callable[[bytes, str], str]):
    """
    设置图片识别回调函数
    callback: 接收 (image_bytes, mime_type) 返回识别结果文本
    """
    global _image_recognition_callback
    _image_recognition_callback = callback


def recognize_image(image_bytes: bytes, mime_type: str = "image/png") -> Optional[str]:
    """调用图片识别回调"""
    if _image_recognition_callback:
        try:
            return _image_recognition_callback(image_bytes, mime_type)
        except Exception:
            return None
    return None


def extract_text_from_file(file_path: str, extract_images: bool = False) -> str:
    """
    根据文件扩展名自动选择解析方法提取文本
    extract_images: 是否提取并识别文档内嵌图片
    """
    ext = os.path.splitext(file_path)[1].lower()
    
    # 图片格式单独处理
    image_extensions = ['.png', '.jpg', '.jpeg', '.gif', '.bmp', '.webp']
    if ext in image_extensions:
        return extract_image_file(file_path)
    
    extractors = {
        '.pdf': lambda p: extract_pdf(p, extract_images),
        '.docx': lambda p: extract_docx(p, extract_images),
        '.doc': extract_doc,
        '.pptx': lambda p: extract_pptx(p, extract_images),
        '.xlsx': extract_xlsx,
        '.xls': extract_xlsx,
        '.txt': extract_text,
        '.md': extract_text,
        '.csv': extract_text,
        '.json': extract_text,
        '.xml': extract_text,
        '.html': extract_html,
        '.htm': extract_html,
    }
    
    extractor = extractors.get(ext)
    if not extractor:
        raise ValueError(f"不支持的文件格式: {ext}")
    
    return extractor(file_path)


def extract_image_file(file_path: str) -> str:
    """提取单独图片文件的内容"""
    if not _image_recognition_callback:
        raise ValueError("未配置图片识别方案，请在知识库设置中选择图片识别方案")
    
    ext = os.path.splitext(file_path)[1].lower()
    mime_types = {
        '.png': 'image/png',
        '.jpg': 'image/jpeg',
        '.jpeg': 'image/jpeg',
        '.gif': 'image/gif',
        '.bmp': 'image/bmp',
        '.webp': 'image/webp',
    }
    mime_type = mime_types.get(ext, 'image/png')
    
    with open(file_path, 'rb') as f:
        image_bytes = f.read()
    
    result = recognize_image(image_bytes, mime_type)
    if result:
        return f"[图片内容]\n{result}"
    raise ValueError("图片识别失败")


def extract_pdf(file_path: str, extract_images: bool = False) -> str:
    """提取 PDF 文本和图片"""
    try:
        from PyPDF2 import PdfReader
        reader = PdfReader(file_path)
        text_parts = []
        has_text = False
        
        for page_num, page in enumerate(reader.pages, 1):
            page_text = page.extract_text()
            if page_text and page_text.strip():
                text_parts.append(f"[第 {page_num} 页]\n{page_text}")
                has_text = True
        
        # 提取PDF中的图片
        if extract_images and _image_recognition_callback:
            image_texts = extract_pdf_images(file_path)
            if image_texts:
                text_parts.extend(image_texts)
        
        # 如果提取到文本，直接返回
        if text_parts:
            return "\n\n".join(text_parts)
        
        raise ValueError("PDF 中未找到可提取的文本")
        
    except ImportError:
        raise ImportError("请安装 PyPDF2: pip install PyPDF2")
    except Exception as e:
        raise ValueError(f"PDF 解析失败: {e}")


def extract_pdf_images(file_path: str) -> List[str]:
    """从PDF中提取图片并识别"""
    image_texts = []
    try:
        import fitz  # PyMuPDF
        
        doc = fitz.open(file_path)
        image_count = 0
        max_images = 20  # 限制最大图片数量
        
        for page_num in range(len(doc)):
            if image_count >= max_images:
                break
                
            page = doc[page_num]
            image_list = page.get_images()
            
            for img_index, img in enumerate(image_list):
                if image_count >= max_images:
                    break
                    
                try:
                    xref = img[0]
                    base_image = doc.extract_image(xref)
                    image_bytes = base_image["image"]
                    image_ext = base_image["ext"]
                    
                    # 跳过太小的图片（可能是图标）
                    if len(image_bytes) < 5000:
                        continue
                    
                    mime_type = f"image/{image_ext}" if image_ext else "image/png"
                    result = recognize_image(image_bytes, mime_type)
                    
                    if result:
                        image_count += 1
                        image_texts.append(f"[第 {page_num + 1} 页 - 图片 {img_index + 1}]\n{result}")
                except Exception:
                    continue
        
        doc.close()
    except ImportError:
        # PyMuPDF 未安装，跳过图片提取
        pass
    except Exception:
        pass
    
    return image_texts


def extract_docx(file_path: str, extract_images: bool = False) -> str:
    """提取 DOCX 文本和图片"""
    try:
        from docx import Document
        from docx.opc.constants import RELATIONSHIP_TYPE as RT
        import zipfile
        
        doc = Document(file_path)
        text_parts = []
        
        # 提取段落文本
        for para in doc.paragraphs:
            if para.text.strip():
                text_parts.append(para.text)
        
        # 提取表格内容
        for table in doc.tables:
            for row in table.rows:
                row_text = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                if row_text:
                    text_parts.append(" | ".join(row_text))
        
        # 提取图片
        if extract_images and _image_recognition_callback:
            image_texts = extract_docx_images(file_path)
            if image_texts:
                text_parts.extend(image_texts)
        
        return "\n".join(text_parts)
    except ImportError:
        raise ImportError("请安装 python-docx: pip install python-docx")
    except Exception as e:
        raise ValueError(f"DOCX 解析失败: {e}")


def extract_docx_images(file_path: str) -> List[str]:
    """从DOCX中提取图片并识别"""
    image_texts = []
    try:
        import zipfile
        
        with zipfile.ZipFile(file_path, 'r') as zip_ref:
            image_count = 0
            max_images = 20
            
            for file_name in zip_ref.namelist():
                if image_count >= max_images:
                    break
                    
                # Word 图片通常在 word/media/ 目录下
                if file_name.startswith('word/media/'):
                    ext = os.path.splitext(file_name)[1].lower()
                    if ext in ['.png', '.jpg', '.jpeg', '.gif', '.bmp', '.webp', '.emf', '.wmf']:
                        try:
                            image_bytes = zip_ref.read(file_name)
                            
                            # 跳过太小的图片
                            if len(image_bytes) < 5000:
                                continue
                            
                            # EMF/WMF 格式跳过（矢量图）
                            if ext in ['.emf', '.wmf']:
                                continue
                            
                            mime_type = f"image/{ext[1:]}" if ext != '.jpg' else "image/jpeg"
                            result = recognize_image(image_bytes, mime_type)
                            
                            if result:
                                image_count += 1
                                image_texts.append(f"[文档图片 {image_count}]\n{result}")
                        except Exception:
                            continue
    except Exception:
        pass
    
    return image_texts


def extract_doc(file_path: str) -> str:
    """提取 DOC 文本 (旧版 Word 格式)"""
    try:
        return extract_docx(file_path, False)
    except:
        pass
    
    try:
        import subprocess
        result = subprocess.run(['antiword', file_path], capture_output=True, text=True)
        if result.returncode == 0:
            return result.stdout
    except:
        pass
    
    try:
        with open(file_path, 'rb') as f:
            content = f.read()
            import re
            text = content.decode('utf-8', errors='ignore')
            text = re.sub(r'[^\x20-\x7E\u4e00-\u9fff\n]', ' ', text)
            text = re.sub(r'\s+', ' ', text).strip()
            if len(text) > 100:
                return text
    except:
        pass
    
    raise ValueError("DOC 格式解析失败，建议转换为 DOCX 格式后上传")


def extract_pptx(file_path: str, extract_images: bool = False) -> str:
    """提取 PPTX 文本和图片"""
    try:
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE
        import zipfile
        
        prs = Presentation(file_path)
        text_parts = []
        
        for slide_num, slide in enumerate(prs.slides, 1):
            slide_texts = [f"[幻灯片 {slide_num}]"]
            
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_texts.append(shape.text)
                if shape.has_table:
                    for row in shape.table.rows:
                        row_text = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                        if row_text:
                            slide_texts.append(" | ".join(row_text))
            
            if len(slide_texts) > 1:
                text_parts.append("\n".join(slide_texts))
        
        # 提取图片
        if extract_images and _image_recognition_callback:
            image_texts = extract_pptx_images(file_path)
            if image_texts:
                text_parts.extend(image_texts)
        
        return "\n\n".join(text_parts)
    except ImportError:
        raise ImportError("请安装 python-pptx: pip install python-pptx")
    except Exception as e:
        raise ValueError(f"PPTX 解析失败: {e}")


def extract_pptx_images(file_path: str) -> List[str]:
    """从PPTX中提取图片并识别"""
    image_texts = []
    try:
        import zipfile
        
        with zipfile.ZipFile(file_path, 'r') as zip_ref:
            image_count = 0
            max_images = 30  # PPT可能有更多图片
            
            for file_name in zip_ref.namelist():
                if image_count >= max_images:
                    break
                    
                # PPT 图片通常在 ppt/media/ 目录下
                if file_name.startswith('ppt/media/'):
                    ext = os.path.splitext(file_name)[1].lower()
                    if ext in ['.png', '.jpg', '.jpeg', '.gif', '.bmp', '.webp']:
                        try:
                            image_bytes = zip_ref.read(file_name)
                            
                            # 跳过太小的图片
                            if len(image_bytes) < 5000:
                                continue
                            
                            mime_type = f"image/{ext[1:]}" if ext != '.jpg' else "image/jpeg"
                            result = recognize_image(image_bytes, mime_type)
                            
                            if result:
                                image_count += 1
                                image_texts.append(f"[PPT图片 {image_count}]\n{result}")
                        except Exception:
                            continue
    except Exception:
        pass
    
    return image_texts


def extract_xlsx(file_path: str) -> str:
    """提取 XLSX/XLS 文本"""
    try:
        from openpyxl import load_workbook
        wb = load_workbook(file_path, data_only=True)
        text_parts = []
        for sheet_name in wb.sheetnames:
            sheet = wb[sheet_name]
            sheet_texts = [f"[工作表: {sheet_name}]"]
            for row in sheet.iter_rows():
                row_values = []
                for cell in row:
                    if cell.value is not None:
                        row_values.append(str(cell.value))
                if row_values:
                    sheet_texts.append(" | ".join(row_values))
            if len(sheet_texts) > 1:
                text_parts.append("\n".join(sheet_texts))
        return "\n\n".join(text_parts)
    except ImportError:
        raise ImportError("请安装 openpyxl: pip install openpyxl")
    except Exception as e:
        raise ValueError(f"XLSX 解析失败: {e}")


def extract_text(file_path: str) -> str:
    """提取纯文本文件"""
    try:
        import chardet
        with open(file_path, 'rb') as f:
            raw = f.read()
            detected = chardet.detect(raw)
            encoding = detected.get('encoding', 'utf-8') or 'utf-8'
        return raw.decode(encoding)
    except ImportError:
        pass
    
    encodings = ['utf-8', 'gbk', 'gb2312', 'utf-16', 'latin-1']
    for enc in encodings:
        try:
            with open(file_path, 'r', encoding=enc) as f:
                return f.read()
        except UnicodeDecodeError:
            continue
    
    raise ValueError("无法识别文件编码")


def extract_html(file_path: str) -> str:
    """提取 HTML 文本"""
    import re
    content = extract_text(file_path)
    content = re.sub(r'<script[^>]*>.*?</script>', '', content, flags=re.DOTALL | re.IGNORECASE)
    content = re.sub(r'<style[^>]*>.*?</style>', '', content, flags=re.DOTALL | re.IGNORECASE)
    content = re.sub(r'<[^>]+>', ' ', content)
    content = re.sub(r'\s+', ' ', content).strip()
    return content


def get_supported_extensions() -> list:
    """返回支持的文件扩展名列表"""
    return ['.pdf', '.docx', '.doc', '.pptx', '.xlsx', '.xls', '.txt', '.md', '.csv', '.json', '.xml', '.html', '.htm', '.png', '.jpg', '.jpeg', '.gif', '.bmp', '.webp']
